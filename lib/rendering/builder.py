"""
Professional Slide Builder - Layer 4A

Builds executive-ready PowerPoint slides using python-pptx with Analytics template styling.
"""

from typing import Dict, List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lib.analysis.insights import Insight, BulletPoint
from PIL import Image
import io

try:
    from lib.rendering.chart_builder_mpl import render_chart_to_png as _mpl_render
    _MPL_AVAILABLE = True
except ImportError:
    _MPL_AVAILABLE = False

# ---------------------------------------------------------------------------
# 3-row chart slide layout constants (13.333" × 7.5")
# ---------------------------------------------------------------------------
_CHART_LEFT  = Inches(0.40)
_CHART_WIDTH = Inches(4.30)
_TEXT_LEFT   = Inches(5.05)
_TEXT_WIDTH  = Inches(7.93)
_ROW_TOPS    = [Inches(1.50), Inches(3.15), Inches(4.80)]
_ROW_H       = Inches(1.55)


def _normalize_image_orientation(image: Image.Image) -> Image.Image:
    """
    Normalize image orientation before inserting into a slide.

    Handles two common sources of rotated/tilted dashboard images:
      1. EXIF orientation metadata (applied transparently via exif_transpose).
      2. 90-degree tilt artifacts from PDF export pipelines — a landscape
         Power BI dashboard can arrive as a portrait image when the exporting
         tool bakes a 90° clockwise rotation into the page stream.  We detect
         this by checking whether the image is taller than it is wide and
         rotate it 90° clockwise (rotate(-90) in PIL) to restore landscape orientation.

    Text readability (left-to-right) is preserved because Power BI dashboards
    are always 16:9 landscape; portrait dimensions unambiguously mean rotation.
    """
    from PIL import ImageOps

    # Step 1: honour EXIF rotation tag (safe no-op if tag is absent)
    try:
        image = ImageOps.exif_transpose(image)
    except Exception:
        pass

    # Step 2: correct 90-degree tilt — portrait image from a landscape source.
    # Power BI embeds landscape dashboards in portrait PDF pages using a CCW
    # rotation; the inverse (90° CW) restores the correct landscape orientation.
    width, height = image.size
    if height > width:
        image = image.rotate(-90, expand=True)

    return image


class AnalyticsStyleGuide:
    """Color palette and typography from Analytics template"""

    # Colors
    DARK_BLUE = RGBColor(0, 32, 96)
    ACCENT_BLUE = RGBColor(0, 114, 198)
    DARK_GRAY = RGBColor(64, 64, 64)
    BEIGE_BG = RGBColor(245, 239, 231)
    WHITE = RGBColor(255, 255, 255)
    PURPLE_ACCENT = RGBColor(112, 48, 160)  # Purple accent for bottom line

    # Typography
    FONT_NAME = "Segoe UI"
    HEADLINE_SIZE = Pt(24)
    INSIGHT_SIZE = Pt(14)
    LABEL_SIZE = Pt(12)

    # Layout measurements (in inches) - 16:9 widescreen format
    SLIDE_WIDTH = Inches(13.333)  # 16:9 aspect ratio
    SLIDE_HEIGHT = Inches(7.5)

    MARGIN = Inches(0.5)
    IMAGE_MAX_WIDTH = Inches(6.5)  # Wider for 16:9
    IMAGE_MAX_HEIGHT = Inches(5.5)
    TEXT_BOX_WIDTH = Inches(5.5)  # Wider text area
    SPACING = Inches(0.4)


class SlideBuilder:
    """Builds professional PowerPoint slides"""

    def __init__(self, source_prs: Presentation):
        """
        Initialize builder with source presentation.

        Args:
            source_prs: Source PowerPoint presentation to extract images from
        """
        self.source_prs = source_prs
        self.style = AnalyticsStyleGuide()

        # Create new presentation
        self.prs = Presentation()
        self.prs.slide_width = int(self.style.SLIDE_WIDTH)
        self.prs.slide_height = int(self.style.SLIDE_HEIGHT)

    def _add_bottom_accent_line(self, slide):
        """Add purple accent line at bottom of slide"""
        # Line positioned at absolute bottom
        line_height = Pt(3)  # Thin line
        line_top = self.style.SLIDE_HEIGHT - line_height  # Right at bottom edge
        line_left = Inches(0)
        line_width = self.style.SLIDE_WIDTH

        # Add rectangle shape as accent line
        line_shape = slide.shapes.add_shape(
            1,  # Rectangle shape type
            line_left,
            line_top,
            line_width,
            line_height
        )

        # Style the line
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self.style.PURPLE_ACCENT
        line_shape.line.fill.background()  # No border

        return line_shape

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add title slide with footer disclaimer"""
        slide_layout = self.prs.slide_layouts[0]  # Title layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Define consistent positioning for title and subtitle
        title_left = self.style.MARGIN
        title_width = self.style.SLIDE_WIDTH - (2 * self.style.MARGIN)
        title_top = Inches(2.5)
        title_height = Inches(1.5)

        subtitle_top = Inches(4.2)
        subtitle_height = Inches(1.0)

        # Set title with explicit dimensions
        title_shape = slide.shapes.title
        title_shape.left = title_left
        title_shape.top = title_top
        title_shape.width = title_width
        title_shape.height = title_height
        title_shape.text = title
        self._style_text(title_shape.text_frame, font_size=Pt(32), bold=True, color=self.style.DARK_BLUE)

        # Set subtitle if provided (aligned with title, explicit dimensions)
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.left = title_left
            subtitle_shape.top = subtitle_top
            subtitle_shape.width = title_width
            subtitle_shape.height = subtitle_height
            subtitle_shape.text = subtitle
            self._style_text(subtitle_shape.text_frame, font_size=Pt(18), color=self.style.DARK_GRAY)

        # Add footer disclaimer
        footer_left = self.style.MARGIN
        footer_top = self.style.SLIDE_HEIGHT - Inches(1)
        footer_width = self.style.SLIDE_WIDTH - (2 * self.style.MARGIN)
        footer_height = Inches(0.5)

        footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
        footer_frame = footer_box.text_frame
        footer_frame.text = "AI generated. Verify before sharing"
        self._style_text(
            footer_frame,
            font_size=Pt(10),
            color=self.style.DARK_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Add purple accent line at bottom
        self._add_bottom_accent_line(slide)

        return slide

    def add_section_slide(self, section_title: str):
        """Add section divider slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add section title centered
        left = self.style.MARGIN
        top = Inches(3)
        width = self.style.SLIDE_WIDTH - (2 * self.style.MARGIN)
        height = Inches(1.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = section_title

        self._style_text(
            text_frame,
            font_size=Pt(36),
            bold=True,
            color=self.style.ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER
        )

        # Add purple accent line at bottom
        self._add_bottom_accent_line(slide)

        return slide

    def add_executive_summary_slide(self, summary_bullets: List[str]):
        """Add executive summary slide with 5 key insights"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            self.style.MARGIN,
            self.style.MARGIN,
            self.style.SLIDE_WIDTH - (2 * self.style.MARGIN),
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        self._style_text(
            title_frame,
            font_size=Pt(32),
            bold=True,
            color=self.style.DARK_BLUE
        )

        # Add bullets
        bullets_top = Inches(1.8)
        bullets_height = self.style.SLIDE_HEIGHT - bullets_top - self.style.MARGIN

        bullets_box = slide.shapes.add_textbox(
            self.style.MARGIN,
            bullets_top,
            self.style.SLIDE_WIDTH - (2 * self.style.MARGIN),
            bullets_height
        )

        text_frame = bullets_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for i, bullet in enumerate(summary_bullets[:5]):  # Max 5 bullets
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.level = 0
            p.space_after = Pt(18)

            # Add bullet with selective bold formatting
            self._add_insight_with_selective_bold(
                p,
                bullet,
                Pt(16),
                self.style.DARK_GRAY
            )

        # Add purple accent line at bottom
        self._add_bottom_accent_line(slide)

        return slide

    def add_recommendations_slide(self, recommendations: List[str]):
        """Add next steps/recommendations slide with 3-5 actions"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            self.style.MARGIN,
            self.style.MARGIN,
            self.style.SLIDE_WIDTH - (2 * self.style.MARGIN),
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Next Steps & Recommendations"
        self._style_text(
            title_frame,
            font_size=Pt(32),
            bold=True,
            color=self.style.ACCENT_BLUE
        )

        # Add recommendations
        recs_top = Inches(1.8)
        recs_height = self.style.SLIDE_HEIGHT - recs_top - self.style.MARGIN

        recs_box = slide.shapes.add_textbox(
            self.style.MARGIN,
            recs_top,
            self.style.SLIDE_WIDTH - (2 * self.style.MARGIN),
            recs_height
        )

        text_frame = recs_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for i, rec in enumerate(recommendations):  # 3-5 recommendations
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.level = 0
            p.space_after = Pt(18)

            # Add recommendation with selective bold formatting
            self._add_insight_with_selective_bold(
                p,
                rec,
                Pt(16),
                self.style.DARK_GRAY,
                prefix=f"{i+1}. "
            )

        # Add purple accent line at bottom
        self._add_bottom_accent_line(slide)

        return slide

    def add_insight_slide(
        self,
        slide_number: int,
        headline: str,
        insights: List[str],
        source_image: Optional[Image.Image] = None
    ):
        """
        Add content slide with headline, insights, and optional screenshot.

        Args:
            slide_number: Slide number from source
            headline: Insight-driven headline
            insights: List of 2-3 bullet points
            source_image: Optional screenshot from source slide
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Calculate layout positions
        image_left = self.style.MARGIN
        image_top = Inches(1.8)  # Moved down from 1.5 for better spacing

        text_left = image_left + self.style.IMAGE_MAX_WIDTH + self.style.SPACING
        text_width = self.style.SLIDE_WIDTH - text_left - self.style.MARGIN

        # Add headline
        headline_box = slide.shapes.add_textbox(
            self.style.MARGIN,
            self.style.MARGIN,
            self.style.SLIDE_WIDTH - (2 * self.style.MARGIN),
            Inches(1)
        )
        headline_frame = headline_box.text_frame
        headline_frame.text = headline
        headline_frame.word_wrap = True
        self._style_text(
            headline_frame,
            font_size=self.style.HEADLINE_SIZE,
            bold=True,
            color=self.style.ACCENT_BLUE
        )

        # Add source image if provided
        if source_image:
            self._add_image_with_border(
                slide,
                source_image,
                image_left,
                image_top,
                self.style.IMAGE_MAX_WIDTH,
                self.style.IMAGE_MAX_HEIGHT
            )

        # Add insights as bullet points
        insights_top = image_top
        insights_height = self.style.IMAGE_MAX_HEIGHT

        insights_box = slide.shapes.add_textbox(
            text_left,
            insights_top,
            text_width,
            insights_height
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        first_para = True
        for i, insight in enumerate(insights[:3]):  # Max 3 insights
            # Split into bold punchy line and supporting detail using || separator
            parts = [s.strip() for s in insight.split('||', 1)]
            bold_line = parts[0]
            detail = parts[1] if len(parts) > 1 else ""

            # Bold punchy line (6-8 words)
            if first_para:
                p = text_frame.paragraphs[0]
                first_para = False
            else:
                p = text_frame.add_paragraph()
                p.space_before = Pt(12)

            p.text = f"• {bold_line}"
            self._style_paragraph(
                p,
                font_size=self.style.INSIGHT_SIZE,
                bold=True,
                color=self.style.DARK_BLUE,
                space_after=Pt(3)
            )

            # Supporting evidence with data (normal weight, below bold line)
            if detail:
                p2 = text_frame.add_paragraph()
                p2.text = detail
                self._style_paragraph(
                    p2,
                    font_size=Pt(13),
                    color=self.style.DARK_GRAY,
                    space_after=Pt(2)
                )

        # Add purple accent line at bottom
        self._add_bottom_accent_line(slide)

        return slide

    def _add_run_with_number_highlight(self, paragraph, text, font_size, base_color,
                                        accent_color=None, bold=True):
        """Split text into alternating normal/number runs; color numbers in accent_color."""
        import re as _re
        if accent_color is None:
            accent_color = self.style.ACCENT_BLUE
        NUMBER_RE = _re.compile(
            r'\b\d[\d,]*(?:\.\d+)?(?:[KMBkm%]|\s*(?:K|M|B|billion|million|thousand))?\b'
        )
        paragraph.text = ""
        pos = 0
        for m in NUMBER_RE.finditer(text):
            if pos < m.start():
                run = paragraph.add_run()
                run.text = text[pos:m.start()]
                run.font.name  = self.style.FONT_NAME
                run.font.size  = font_size
                run.font.bold  = bold
                run.font.color.rgb = base_color
            run = paragraph.add_run()
            run.text = m.group()
            run.font.name  = self.style.FONT_NAME
            run.font.size  = font_size
            run.font.bold  = bold
            run.font.color.rgb = accent_color
            pos = m.end()
        if pos < len(text):
            run = paragraph.add_run()
            run.text = text[pos:]
            run.font.name  = self.style.FONT_NAME
            run.font.size  = font_size
            run.font.bold  = bold
            run.font.color.rgb = base_color

    def add_polished_chart_slide(
        self,
        slide_number: int,
        headline: str,
        bullet_points: List[BulletPoint],
        use_mpl: bool = True,
    ):
        """
        'Work of art' layout — branches on number of charts:

        n_charts == 1  (Commentary Layout):
          Headline full-width (Pt 20)
          Chart LEFT 7.10" × 6.10"  |  Insight panel RIGHT 5.48"
          Each insight block: left accent bar + bold Pt(16) + detail Pt(12)

        n_charts == 2  (Expanded Strip Layout):
          Headline full-width (Pt 20)
          Side-by-side charts 3.65" tall
          Expanded 3-column insight strip 2.10" at Pt(14)/Pt(11)

        n_charts == 0  (fallback):
          Insight strip only (same as original).
        """
        from lib.rendering.chart_builder import render_chart

        MARGIN   = Inches(0.25)
        USABLE_W = self.style.SLIDE_WIDTH - 2 * MARGIN  # 12.833"

        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # ── Headline (Pt 20, taller box) ────────────────────────────────────
        hl_box = slide.shapes.add_textbox(MARGIN, Inches(0.10), USABLE_W, Inches(0.65))
        hl_tf  = hl_box.text_frame
        hl_tf.word_wrap = True
        hl_tf.paragraphs[0].text = headline
        self._style_text(hl_tf, font_size=Pt(20), bold=True,
                         color=self.style.DARK_BLUE)

        # ── Identify chart bullet-points (max 2) ────────────────────────────
        chart_bps = [bp for bp in bullet_points
                     if getattr(bp, 'chart', None) is not None][:2]
        n_charts  = len(chart_bps)

        # ════════════════════════════════════════════════════════════════════
        # Branch A: 1 chart — Commentary Layout
        # ════════════════════════════════════════════════════════════════════
        if n_charts == 1:
            CONTENT_TOP  = Inches(0.85)
            CONTENT_H    = Inches(6.10)
            CHART_W      = Inches(7.10)
            GAP          = Inches(0.25)
            INS_LEFT     = MARGIN + CHART_W + GAP        # 7.60"
            INS_W        = Inches(5.48)
            BLOCK_H      = Inches(1.90)                  # 3 × 1.90 = 5.70 fits in 6.10"
            ACCENT_BAR_W = Inches(0.05)
            TEXT_INDENT  = Inches(0.12)
            TEXT_W       = INS_W - ACCENT_BAR_W - TEXT_INDENT - Inches(0.05)

            spec = chart_bps[0].chart
            if use_mpl and _MPL_AVAILABLE:
                w_in = CHART_W   / 914400
                h_in = CONTENT_H / 914400
                png  = _mpl_render(spec, w_in, h_in, dpi=200)
                slide.shapes.add_picture(io.BytesIO(png), MARGIN, CONTENT_TOP,
                                         CHART_W, CONTENT_H)
            else:
                title_text = (spec.title or "").strip()
                if title_text:
                    tb = slide.shapes.add_textbox(MARGIN, CONTENT_TOP - Inches(0.22),
                                                  CHART_W, Inches(0.22))
                    tf = tb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    p.text = title_text
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name  = self.style.FONT_NAME
                        run.font.size  = Pt(10)
                        run.font.bold  = False
                        run.font.color.rgb = RGBColor(89, 89, 89)
                render_chart(slide, MARGIN, CONTENT_TOP, CHART_W, CONTENT_H, spec)

            # Insight blocks stacked vertically on the right
            for i, bp in enumerate(bullet_points[:3]):
                block_top = CONTENT_TOP + i * BLOCK_H

                # Left accent bar (ACCENT_BLUE, ~3pt wide)
                bar = slide.shapes.add_shape(
                    1,
                    INS_LEFT, block_top + Inches(0.05),
                    ACCENT_BAR_W, BLOCK_H - Inches(0.10)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = self.style.ACCENT_BLUE
                bar.line.fill.background()

                # Text box
                tb = slide.shapes.add_textbox(
                    INS_LEFT + ACCENT_BAR_W + TEXT_INDENT, block_top,
                    TEXT_W, BLOCK_H
                )
                tf = tb.text_frame
                tf.word_wrap = True

                parts     = [s.strip() for s in bp.text.split('||', 1)]
                bold_line = parts[0]
                detail    = parts[1] if len(parts) > 1 else ""

                # Bold line with number highlighting (Pt 16)
                p = tf.paragraphs[0]
                self._add_run_with_number_highlight(
                    p, f"\u2022 {bold_line}", Pt(16), self.style.DARK_BLUE
                )
                p.space_after = Pt(5)

                # Detail (Pt 12, DARK_GRAY)
                if detail:
                    p2 = tf.add_paragraph()
                    p2.text = detail
                    self._style_paragraph(p2, font_size=Pt(12),
                                          color=self.style.DARK_GRAY)

        # ════════════════════════════════════════════════════════════════════
        # Branch B: 2 charts — Expanded Strip Layout
        # ════════════════════════════════════════════════════════════════════
        elif n_charts == 2:
            CTITLE_Y = Inches(0.78)   # after headline (ends at 0.75")
            CTITLE_H = Inches(0.18)
            CHART_Y  = Inches(0.96)
            CHART_H  = Inches(3.65)   # was 4.90

            gap = Inches(0.15)
            cw  = (USABLE_W - gap) // 2
            positions = [
                (MARGIN,            cw),
                (MARGIN + cw + gap, cw),
            ]

            for i, (cx, cw) in enumerate(positions):
                spec = chart_bps[i].chart
                if use_mpl and _MPL_AVAILABLE:
                    w_in = cw     / 914400
                    h_in = CHART_H / 914400
                    png  = _mpl_render(spec, w_in, h_in, dpi=200)
                    slide.shapes.add_picture(io.BytesIO(png), cx, CHART_Y, cw, CHART_H)
                else:
                    title_text = (spec.title or "").strip()
                    if title_text:
                        tb = slide.shapes.add_textbox(cx, CTITLE_Y, cw, CTITLE_H)
                        tf = tb.text_frame
                        tf.word_wrap = False
                        p = tf.paragraphs[0]
                        p.text = title_text
                        p.alignment = PP_ALIGN.LEFT
                        for run in p.runs:
                            run.font.name  = self.style.FONT_NAME
                            run.font.size  = Pt(10)
                            run.font.bold  = False
                            run.font.color.rgb = RGBColor(89, 89, 89)
                    render_chart(slide, cx, CHART_Y, cw, CHART_H, spec)

            # Thin separator
            sep_y = CHART_Y + CHART_H + Inches(0.03)
            sep   = slide.shapes.add_shape(1, MARGIN, sep_y, USABLE_W, Pt(0.75))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(220, 220, 220)
            sep.line.fill.background()

            # Expanded 3-column insight strip
            INS_Y   = sep_y + Pt(3)
            INS_H   = Inches(2.10)   # was 1.15
            col_w   = USABLE_W // 3
            col_pad = Inches(0.08)

            for i, bp in enumerate(bullet_points[:3]):
                col_left = MARGIN + i * col_w

                if i > 0:
                    vsep_x = col_left - Pt(0.5)
                    vsep   = slide.shapes.add_shape(
                        1, vsep_x, INS_Y, Pt(0.75), INS_H)
                    vsep.fill.solid()
                    vsep.fill.fore_color.rgb = RGBColor(220, 220, 220)
                    vsep.line.fill.background()

                tb = slide.shapes.add_textbox(
                    col_left + col_pad, INS_Y,
                    col_w - 2 * col_pad, INS_H
                )
                tf = tb.text_frame
                tf.word_wrap = True

                parts     = [s.strip() for s in bp.text.split('||', 1)]
                bold_line = parts[0]
                detail    = parts[1] if len(parts) > 1 else ""

                # Bold line with number highlighting (Pt 14)
                p = tf.paragraphs[0]
                self._add_run_with_number_highlight(
                    p, f"\u2022 {bold_line}", Pt(14), self.style.DARK_BLUE
                )
                p.space_after = Pt(2)

                if detail:
                    p2 = tf.add_paragraph()
                    p2.text = detail
                    self._style_paragraph(p2, font_size=Pt(11),
                                          color=self.style.DARK_GRAY)

        # ════════════════════════════════════════════════════════════════════
        # Branch C: 0 charts — fallback insight strip (original behaviour)
        # ════════════════════════════════════════════════════════════════════
        else:
            sep_y = Inches(0.96) + Inches(4.90) + Inches(0.03)
            sep   = slide.shapes.add_shape(1, MARGIN, sep_y, USABLE_W, Pt(0.75))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(220, 220, 220)
            sep.line.fill.background()

            INS_Y   = sep_y + Pt(3)
            INS_H   = Inches(1.15)
            col_w   = USABLE_W // 3
            col_pad = Inches(0.08)

            for i, bp in enumerate(bullet_points[:3]):
                col_left = MARGIN + i * col_w

                if i > 0:
                    vsep_x = col_left - Pt(0.5)
                    vsep   = slide.shapes.add_shape(
                        1, vsep_x, INS_Y, Pt(0.75), INS_H)
                    vsep.fill.solid()
                    vsep.fill.fore_color.rgb = RGBColor(220, 220, 220)
                    vsep.line.fill.background()

                tb = slide.shapes.add_textbox(
                    col_left + col_pad, INS_Y,
                    col_w - 2 * col_pad, INS_H
                )
                tf = tb.text_frame
                tf.word_wrap = True

                parts     = [s.strip() for s in bp.text.split('||', 1)]
                bold_line = parts[0]
                detail    = parts[1] if len(parts) > 1 else ""

                p = tf.paragraphs[0]
                p.text = f"\u2022 {bold_line}"
                self._style_paragraph(p, font_size=Pt(11), bold=True,
                                      color=self.style.DARK_BLUE, space_after=Pt(2))
                if detail:
                    p2 = tf.add_paragraph()
                    p2.text = detail
                    self._style_paragraph(p2, font_size=Pt(9),
                                          color=self.style.DARK_GRAY)

        # ── Footer ──────────────────────────────────────────────────────────
        footer_box = slide.shapes.add_textbox(
            MARGIN, Inches(7.08), USABLE_W, Inches(0.25))
        ff = footer_box.text_frame
        ff.paragraphs[0].text = "AI generated. Verify before sharing"
        self._style_text(ff, font_size=Pt(8), color=self.style.DARK_GRAY,
                         alignment=PP_ALIGN.CENTER)

        self._add_bottom_accent_line(slide)
        return slide

    def add_chart_insight_slide(
        self,
        slide_number: int,
        headline: str,
        bullet_points: List[BulletPoint],
        section_tag: str = ""
    ):
        """
        3-row slide layout: each row = [mini native chart left | insight text right].

        Falls back to add_insight_slide (image-left) if all charts are None.
        """
        from lib.rendering.chart_builder import render_chart

        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Section tag
        if section_tag:
            tag_box = slide.shapes.add_textbox(
                self.style.MARGIN, Inches(0.25),
                self.style.SLIDE_WIDTH - 2 * self.style.MARGIN, Inches(0.45)
            )
            tf = tag_box.text_frame
            tf.paragraphs[0].text = section_tag.upper()
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.font.name = self.style.FONT_NAME
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = self.style.ACCENT_BLUE

        # Headline
        headline_box = slide.shapes.add_textbox(
            self.style.MARGIN, Inches(0.75),
            self.style.SLIDE_WIDTH - 2 * self.style.MARGIN, Inches(0.65)
        )
        hf = headline_box.text_frame
        hf.word_wrap = True
        hf.paragraphs[0].text = headline
        self._style_text(hf, font_size=Pt(20), bold=True, color=self.style.ACCENT_BLUE)

        # Thin ACCENT_BLUE divider at y=1.40"
        divider = slide.shapes.add_shape(
            1,  # rectangle
            self.style.MARGIN, Inches(1.40),
            self.style.SLIDE_WIDTH - 2 * self.style.MARGIN, Pt(1)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.style.ACCENT_BLUE
        divider.line.fill.background()

        # 3 rows
        for row_idx, bp in enumerate(bullet_points[:3]):
            row_top = _ROW_TOPS[row_idx]

            # Chart on left
            if bp.chart is not None:
                render_chart(slide, _CHART_LEFT, row_top, _CHART_WIDTH, _ROW_H, bp.chart)

            # Text on right
            text_box = slide.shapes.add_textbox(
                _TEXT_LEFT, row_top,
                _TEXT_WIDTH, _ROW_H
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            parts = [s.strip() for s in bp.text.split('||', 1)]
            bold_line = parts[0]
            detail = parts[1] if len(parts) > 1 else ""

            # Bold line
            p = tf.paragraphs[0]
            p.text = f"• {bold_line}"
            self._style_paragraph(p, font_size=Pt(13), bold=True,
                                  color=self.style.DARK_BLUE, space_after=Pt(3))

            # Detail line
            if detail:
                p2 = tf.add_paragraph()
                p2.text = detail
                self._style_paragraph(p2, font_size=Pt(11), color=self.style.DARK_GRAY)

            # Row separator (except after last row)
            if row_idx < min(2, len(bullet_points) - 1):
                sep_y = row_top + _ROW_H - Pt(0.5)
                sep = slide.shapes.add_shape(
                    1,
                    self.style.MARGIN, sep_y,
                    self.style.SLIDE_WIDTH - 2 * self.style.MARGIN, Pt(0.5)
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = RGBColor(200, 200, 200)
                sep.line.fill.background()

        # Footer
        footer_box = slide.shapes.add_textbox(
            self.style.MARGIN, Inches(6.90),
            self.style.SLIDE_WIDTH - 2 * self.style.MARGIN, Inches(0.45)
        )
        ff = footer_box.text_frame
        ff.paragraphs[0].text = "AI generated. Verify before sharing"
        self._style_text(ff, font_size=Pt(9), color=self.style.DARK_GRAY,
                         alignment=PP_ALIGN.CENTER)

        self._add_bottom_accent_line(slide)
        return slide

    def _add_image_with_border(
        self,
        slide,
        image: Image.Image,
        left: float,
        top: float,
        max_width: float,
        max_height: float
    ):
        """Add image with white background and blue border"""
        # Normalize orientation before placement: correct EXIF rotation and any
        # 90-degree tilt artifacts introduced by PDF/image export pipelines.
        image = _normalize_image_orientation(image)

        # Calculate scaled dimensions maintaining aspect ratio
        img_width, img_height = image.size
        aspect_ratio = img_width / img_height

        if aspect_ratio > 1:  # Landscape
            width = min(max_width, Inches(img_width / 96))  # Assume 96 DPI
            height = width / aspect_ratio
        else:  # Portrait
            height = min(max_height, Inches(img_height / 96))
            width = height * aspect_ratio

        # Ensure doesn't exceed max dimensions
        if width > max_width:
            width = max_width
            height = width / aspect_ratio

        if height > max_height:
            height = max_height
            width = height * aspect_ratio

        # Save image to BytesIO
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG')
        img_stream.seek(0)

        # Add image
        pic = slide.shapes.add_picture(img_stream, left, top, width, height)

        # Add border (blue line) - lighter weight
        line = pic.line
        line.color.rgb = self.style.ACCENT_BLUE
        line.width = Pt(1)

        return pic

    def _add_insight_with_selective_bold(self, paragraph, insight_text, font_size, color, prefix="• "):
        """Add insight text with selective bolding for emphasis (supports markdown **bold**)"""
        import re

        # Clear existing text
        paragraph.text = ""

        # Add prefix (bullet or number)
        prefix_run = paragraph.add_run()
        prefix_run.text = prefix
        prefix_run.font.name = self.style.FONT_NAME
        prefix_run.font.size = font_size
        prefix_run.font.color.rgb = color

        # Check for markdown bold syntax (**text**)
        markdown_bold_pattern = r'\*\*(.+?)\*\*'
        markdown_matches = list(re.finditer(markdown_bold_pattern, insight_text))

        if markdown_matches:
            # Use markdown bold sections
            pos = 0
            for match in markdown_matches:
                # Add normal text before bold section
                if pos < match.start():
                    run = paragraph.add_run()
                    run.text = insight_text[pos:match.start()]
                    run.font.name = self.style.FONT_NAME
                    run.font.size = font_size
                    run.font.color.rgb = color

                # Add bold section (without ** markers)
                run = paragraph.add_run()
                run.text = match.group(1)  # Extract text without **
                run.font.name = self.style.FONT_NAME
                run.font.size = font_size
                run.font.color.rgb = color
                run.font.bold = True

                pos = match.end()

            # Add remaining text
            if pos < len(insight_text):
                run = paragraph.add_run()
                run.text = insight_text[pos:]
                run.font.name = self.style.FONT_NAME
                run.font.size = font_size
                run.font.color.rgb = color
        else:
            # No markdown - just add plain text without any bolding
            run = paragraph.add_run()
            run.text = insight_text
            run.font.name = self.style.FONT_NAME
            run.font.size = font_size
            run.font.color.rgb = color

    def _style_text(
        self,
        text_frame,
        font_size: Pt,
        bold: bool = False,
        color: RGBColor = None,
        alignment: PP_ALIGN = PP_ALIGN.LEFT
    ):
        """Apply text styling"""
        for paragraph in text_frame.paragraphs:
            self._style_paragraph(paragraph, font_size, bold, color, alignment)

    def _style_paragraph(
        self,
        paragraph,
        font_size: Pt,
        bold: bool = False,
        color: RGBColor = None,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        space_after: Pt = None
    ):
        """Apply paragraph styling"""
        paragraph.alignment = alignment

        if space_after:
            paragraph.space_after = space_after

        for run in paragraph.runs:
            font = run.font
            font.name = self.style.FONT_NAME
            font.size = font_size
            font.bold = bold

            if color:
                font.color.rgb = color

    def save(self, output_path: str):
        """Save presentation to file"""
        self.prs.save(output_path)


def _get_source_images_from_temp(source_path: str) -> Dict[int, str]:
    """
    Get already-extracted images from temp/ directory.
    Used when source is PDF (images already extracted during prepare phase).

    Args:
        source_path: Path to source file (used for validation)

    Returns:
        Dictionary mapping slide_number → image_path
    """
    import json

    try:
        with open('temp/analysis_request.json', 'r', encoding='utf-8') as f:
            request = json.load(f)

        # Build mapping of slide number to image path
        image_map = {}
        for slide_info in request.get('slides', []):
            slide_num = slide_info['slide_number']
            image_path = slide_info['image_path']
            image_map[slide_num] = image_path

        return image_map

    except (FileNotFoundError, json.JSONDecodeError, KeyError) as e:
        print(f"  WARNING: Could not load temp images: {e}")
        return {}


def extract_slide_image(source_prs: Presentation, slide_number: int) -> Optional[Image.Image]:
    """
    Extract screenshot from source slide.

    Args:
        source_prs: Source presentation
        slide_number: Slide number (0-indexed)

    Returns:
        PIL Image or None if extraction fails
    """
    try:
        if slide_number >= len(source_prs.slides):
            return None

        slide = source_prs.slides[slide_number]

        # Look for images in slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image_stream = io.BytesIO(shape.image.blob)
                return Image.open(image_stream)

        return None

    except Exception:
        return None


def render_presentation(
    source_path: str,
    insights: Dict[str, Insight],
    output_path: str
):
    """
    Main entry point: Render executive presentation.
    Supports both .pptx and .pdf source files.

    Args:
        source_path: Path to source file (.pptx or .pdf)
        insights: Dictionary of slide titles to Insights
        output_path: Path for output PowerPoint
    """
    file_type = Path(source_path).suffix.lower()

    # For PDF and PBIP sources, images are already in temp/ (or absent for PBIP).
    # For PPTX sources, we load the presentation directly.
    if file_type in ('.pdf', '.pbip', '.pbix') or Path(source_path).is_dir():
        # Create a blank dummy presentation (not used for image extraction)
        from pptx import Presentation as PrsClass
        source_prs = PrsClass()
        source_images_map = _get_source_images_from_temp(source_path)
    else:
        # Load source presentation for PPTX
        source_prs = Presentation(source_path)
        source_images_map = None

    # Create builder
    builder = SlideBuilder(source_prs)

    # Add title slide (use deck_title from insights if provided)
    deck_title = insights.get('__deck_title__', 'Executive Insights')
    deck_subtitle = insights.get('__deck_subtitle__', 'Insights & Recommendations')
    builder.add_title_slide(deck_title, deck_subtitle)

    # Add executive summary slide (if provided)
    if '__executive_summary__' in insights:
        executive_summary = insights['__executive_summary__']
        if isinstance(executive_summary, list) and executive_summary:
            builder.add_executive_summary_slide(executive_summary)

    # Add content slides
    slide_mapping = {}  # Track which source slides we've processed

    # For PDF sources, iterate through insights directly (no source slides to iterate)
    # For PPTX sources, iterate through source slides as before
    if source_images_map is not None:
        # PDF workflow: Use temp/ images and insights mapping
        # Load analysis request once (not in loop)
        import json
        slide_info_map = {}
        try:
            with open('temp/analysis_request.json', 'r', encoding='utf-8') as f:
                request = json.load(f)
                for slide_info in request.get('slides', []):
                    slide_info_map[slide_info['slide_number']] = slide_info
        except Exception as e:
            print(f"  WARNING: Could not load analysis request: {e}")

        # Process insights - iterate by slide number for guaranteed order
        for slide_num in sorted(source_images_map.keys()):
            slide_info = slide_info_map.get(slide_num)
            if not slide_info:
                continue

            # Try to find matching insight by slide_number first (preferred)
            insight = insights.get(slide_num)

            # Fallback to title matching for backward compatibility
            if not insight:
                title = slide_info['title']
                insight = insights.get(title)

                if not insight:
                    # Try without special characters if exact match fails
                    import re
                    title_clean = re.sub(r'[^\w\s]', '', title).strip()
                    for key, val in insights.items():
                        if isinstance(key, str) and re.sub(r'[^\w\s]', '', key).strip() == title_clean:
                            insight = val
                            break

            if insight:
                # Load source image
                source_image = None
                image_path = source_images_map.get(slide_num)
                if image_path and Path(image_path).exists():
                    try:
                        source_image = _normalize_image_orientation(Image.open(image_path))
                    except Exception as e:
                        print(f"  WARNING: Could not load image for slide {slide_num}: {e}")

                # Determine whether to use chart-row layout or image-left layout
                has_charts = any(
                    hasattr(bp, 'chart') and bp.chart is not None
                    for bp in insight.bullet_points
                )
                if has_charts:
                    builder.add_polished_chart_slide(
                        slide_number=slide_num,
                        headline=insight.headline,
                        bullet_points=insight.bullet_points
                    )
                else:
                    plain = [
                        bp.text if hasattr(bp, 'text') else str(bp)
                        for bp in insight.bullet_points
                    ]
                    builder.add_insight_slide(
                        slide_number=slide_num,
                        headline=insight.headline,
                        insights=plain,
                        source_image=source_image
                    )
            else:
                print(f"  WARNING: No insight found for slide {slide_num}: {title}")
    else:
        # PPTX workflow: Original logic
        for slide_idx, slide in enumerate(source_prs.slides):
            # Extract slide title (first shape with text, typically)
            slide_title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_title = shape.text.strip()
                    break

            # Remove emoji characters to match parsed titles
            import re
            slide_title_clean = re.sub(r'[\U00010000-\U0010ffff]', '', slide_title).strip()

            # Find matching insight — try slide_number key first (integer),
            # then fall back to title string for backward compatibility
            insight = insights.get(slide_idx + 1) or insights.get(slide_title_clean)

            if insight:
                # Extract image from source slide and normalize orientation
                source_image = extract_slide_image(source_prs, slide_idx)
                if source_image:
                    source_image = _normalize_image_orientation(source_image)

                # Determine whether to use chart-row layout or image-left layout
                has_charts = any(
                    hasattr(bp, 'chart') and bp.chart is not None
                    for bp in insight.bullet_points
                )
                if has_charts:
                    builder.add_polished_chart_slide(
                        slide_number=slide_idx + 1,
                        headline=insight.headline,
                        bullet_points=insight.bullet_points
                    )
                else:
                    plain = [
                        bp.text if hasattr(bp, 'text') else str(bp)
                        for bp in insight.bullet_points
                    ]
                    builder.add_insight_slide(
                        slide_number=slide_idx + 1,
                        headline=insight.headline,
                        insights=plain,
                        source_image=source_image
                    )

    # Add recommendations slide (if provided)
    if '__recommendations__' in insights:
        recommendations = insights['__recommendations__']
        if isinstance(recommendations, list) and recommendations:
            builder.add_recommendations_slide(recommendations)

    # Save presentation
    builder.save(output_path)
