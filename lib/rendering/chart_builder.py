"""
Chart Builder - Native PowerPoint chart rendering.

Polished, Power BI-inspired chart styling:
  • No chart borders (chart floats cleanly on slide)
  • Platform-matched color palettes (purple/Agents, pink/Chat, blue/M365)
  • 9-10pt Segoe UI fonts throughout
  • Subtle light gridlines on value axis
  • Data labels at natural positions
  • Legend when multiple series, positioned cleanly

All render functions signature:
    render_X(slide, left, top, width, height, spec) -> Optional[shape]
"""

from typing import Optional, List
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData, BubbleChartData, XyChartData
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

from lib.analysis.insights import ChartSpec

# ---------------------------------------------------------------------------
# Power BI–matching color palettes
# ---------------------------------------------------------------------------

# Agents (purple family)
PBI_PURPLE_DARK  = RGBColor(91,  31,  110)
PBI_PURPLE       = RGBColor(123, 63,  139)
PBI_PURPLE_MID   = RGBColor(155, 100, 178)
PBI_PURPLE_LIGHT = RGBColor(200, 170, 218)
PBI_PURPLE_PALE  = RGBColor(227, 212, 238)

# Chat – Unlicensed (pink/magenta family)
PBI_PINK_DARK    = RGBColor(160, 48,  112)
PBI_PINK         = RGBColor(200, 75,  151)
PBI_PINK_MID     = RGBColor(224, 122, 187)
PBI_PINK_LIGHT   = RGBColor(240, 185, 218)
PBI_PINK_PALE    = RGBColor(250, 224, 239)

# M365 Copilot (blue family)
PBI_BLUE_DARK    = RGBColor(0,   50,  120)
PBI_BLUE         = RGBColor(0,   114, 198)
PBI_BLUE_MID     = RGBColor(100, 165, 220)
PBI_BLUE_LIGHT   = RGBColor(189, 220, 243)
PBI_BLUE_PALE    = RGBColor(222, 240, 252)

# Neutrals
DARK_GRAY        = RGBColor(89,  89,  89)
MID_GRAY         = RGBColor(140, 140, 140)
LIGHT_GRAY       = RGBColor(200, 200, 200)
GRID_GRAY        = RGBColor(232, 232, 232)
WHITE            = RGBColor(255, 255, 255)
NEAR_BLACK       = RGBColor(38,  38,  38)

# Legacy aliases kept for backward compat
ACCENT_BLUE = PBI_BLUE
DARK_BLUE   = PBI_BLUE_DARK
PURPLE      = PBI_PURPLE

# Multi-series palette (cycles for line/area/scatter)
MULTI_PALETTE = [PBI_BLUE, PBI_PINK, PBI_PURPLE, PBI_BLUE_DARK,
                 PBI_PINK_DARK, PBI_PURPLE_MID, PBI_BLUE_MID]

# Habit-tier donut palettes  (Light → Moderate → Frequent → Daily)
HABIT_AGENTS = [RGBColor(195,195,195), PBI_PURPLE_LIGHT, PBI_PURPLE,      PBI_PURPLE_DARK]
HABIT_CHAT   = [RGBColor(195,195,195), PBI_PINK_LIGHT,   PBI_PINK,        PBI_PINK_DARK]
HABIT_M365   = [RGBColor(195,195,195), PBI_BLUE_LIGHT,   PBI_BLUE,        PBI_BLUE_DARK]
DONUT_PALETTE= MULTI_PALETTE  # general fallback

FONT_NAME = "Segoe UI"


# ---------------------------------------------------------------------------
# Border & style helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(s: str) -> Optional[RGBColor]:
    if not s:
        return None
    s = s.lstrip('#')
    if len(s) != 6:
        return None
    try:
        return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
    except ValueError:
        return None


def _apply_light_chart_border(chart_frame):
    """
    Apply a subtle 0.5 pt light-gray border to the outer chart frame.
    Also strips borders from chart-space and plot-area so the internal
    chart floats cleanly while the outer frame gives it visual structure.
    """
    # 1. Outer graphic-frame: thin light-gray line
    try:
        chart_frame.line.color.rgb = RGBColor(210, 210, 210)
        chart_frame.line.width     = Pt(0.5)
    except Exception:
        pass

    def _set_no_ln(parent_el, ns_prefix='c'):
        """Strip inner <a:ln> from <{ns}:spPr> so internal borders are invisible."""
        spPr_tag = qn(f'{ns_prefix}:spPr')
        spPr = parent_el.find(spPr_tag)
        if spPr is None:
            spPr = etree.SubElement(parent_el, spPr_tag)
        for ln in spPr.findall(qn('a:ln')):
            spPr.remove(ln)
        ln_el = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln_el, qn('a:noFill'))

    try:
        chart_obj  = chart_frame.chart
        chart_space = chart_obj._element.getparent()
        _set_no_ln(chart_space, 'c')
    except Exception:
        pass

    try:
        plot_area = chart_obj._element.find(qn('c:plotArea'))
        if plot_area is not None:
            _set_no_ln(plot_area, 'c')
    except Exception:
        pass


def _enable_data_labels(dl_obj, show_value: bool = True,
                         show_percent: bool = False,
                         show_series_name: bool = False,
                         position: str = None):
    """
    Force-enable data-label display flags via direct XML editing.
    position: OpenXML dLblPos value e.g. 'outEnd', 'inEnd', 'ctr', 'b', 't'
    """
    try:
        el = dl_obj._element

        def _sv(tag, val):
            ch = el.find(qn(tag))
            if ch is None:
                ch = etree.SubElement(el, qn(tag))
            ch.set('val', '1' if val else '0')

        _sv('c:showLegendKey', False)
        _sv('c:showVal',       show_value)
        _sv('c:showCatName',   False)
        _sv('c:showSerName',   show_series_name)
        _sv('c:showPercent',   show_percent)

        if position:
            ch = el.find(qn('c:dLblPos'))
            if ch is None:
                ch = etree.SubElement(el, qn('c:dLblPos'))
            ch.set('val', position)
    except Exception:
        pass


def _set_axis_title(axis, title_text: str, font_size_pt: int = 8):
    """Inject an axis title element via raw XML (no high-level API in python-pptx)."""
    try:
        ax_el = axis._element
        for t in ax_el.findall(qn('c:title')):
            ax_el.remove(t)

        title_el = etree.SubElement(ax_el, qn('c:title'))
        tx_el    = etree.SubElement(title_el, qn('c:tx'))
        rich_el  = etree.SubElement(tx_el, qn('c:rich'))
        etree.SubElement(rich_el, qn('a:bodyPr'))
        etree.SubElement(rich_el, qn('a:lstStyle'))
        p_el     = etree.SubElement(rich_el, qn('a:p'))
        r_el     = etree.SubElement(p_el, qn('a:r'))
        rPr      = etree.SubElement(r_el, qn('a:rPr'))
        rPr.set('lang', 'en-US')
        rPr.set('sz',   str(font_size_pt * 100))
        rPr.set('b',    '0')
        sf  = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(sf,  qn('a:srgbClr'))
        clr.set('val', '595959')
        t_el      = etree.SubElement(r_el, qn('a:t'))
        t_el.text = title_text
        etree.SubElement(title_el, qn('c:overlay')).set('val', '0')
    except Exception:
        pass


def _apply_pbi_chart_style(chart_obj, chart_frame=None, show_gridlines: bool = True):
    """
    Apply polished Power BI–style formatting:
      • Remove chart borders
      • Light gray subtle gridlines on value axis
      • 9pt Segoe UI category labels
      • Hide value-axis labels (data labels carry the numbers)
      • No embedded chart title
      • No legend (callers add it explicitly when needed)
    """
    if chart_frame is not None:
        _apply_light_chart_border(chart_frame)

    # Category (horizontal/X) axis
    try:
        cat_ax = chart_obj.category_axis
        cat_ax.has_major_gridlines = False
        cat_ax.has_minor_gridlines = False
        cat_ax.tick_labels.font.size  = Pt(9)
        cat_ax.tick_labels.font.name  = FONT_NAME
        cat_ax.tick_labels.font.color.rgb = DARK_GRAY
        cat_ax.tick_labels.font.bold  = False
        # Thin axis line
        try:
            cat_ax.format.line.color.rgb  = GRID_GRAY
            cat_ax.format.line.width      = Pt(0.5)
        except Exception:
            pass
    except Exception:
        pass

    # Value (vertical/Y) axis — keep gridlines, hide labels
    try:
        val_ax = chart_obj.value_axis
        val_ax.has_major_gridlines = show_gridlines
        val_ax.has_minor_gridlines = False
        if show_gridlines:
            try:
                val_ax.major_gridlines.format.line.color.rgb = GRID_GRAY
                val_ax.major_gridlines.format.line.width     = Pt(0.5)
            except Exception:
                pass
        # Make axis labels invisible rather than deleting axis (deletion removes gridlines)
        val_ax.tick_labels.font.size  = Pt(1)
        val_ax.tick_labels.font.color.rgb = WHITE
        try:
            val_ax.format.line.fill.background()   # hide axis line
        except Exception:
            pass
    except Exception:
        pass

    # No embedded title
    try:
        chart_obj.chart_title.has_text_frame = False
    except Exception:
        pass

    # No legend by default (add explicitly via _add_legend when multi-series)
    try:
        chart_obj.has_legend = False
    except Exception:
        pass


def _add_legend(chart_obj, position=XL_LEGEND_POSITION.BOTTOM, font_size: int = 9):
    try:
        chart_obj.has_legend = True
        chart_obj.legend.position = position
        chart_obj.legend.include_in_layout = False
        chart_obj.legend.font.size  = Pt(font_size)
        chart_obj.legend.font.name  = FONT_NAME
        chart_obj.legend.font.color.rgb = DARK_GRAY
    except Exception:
        pass


def _set_series_color(series, color: RGBColor):
    try:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
    except Exception:
        pass


def _set_data_labels(series, font_size_pt: int = 9,
                     color: RGBColor = DARK_GRAY, bold: bool = False,
                     show_value: bool = True, show_percent: bool = False,
                     show_series_name: bool = False, position: str = None):
    try:
        dl = series.data_labels
        dl.font.size      = Pt(font_size_pt)
        dl.font.name      = FONT_NAME
        dl.font.color.rgb = color
        dl.font.bold      = bold
        _enable_data_labels(dl, show_value=show_value,
                            show_percent=show_percent,
                            show_series_name=show_series_name,
                            position=position)
    except Exception:
        pass


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size_pt: int = 9, bold: bool = False,
                 color: RGBColor = DARK_GRAY, align=PP_ALIGN.LEFT,
                 word_wrap: bool = True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.name       = FONT_NAME
    run.font.size       = Pt(font_size_pt)
    run.font.bold       = bold
    run.font.color.rgb  = color
    return txb


# ---------------------------------------------------------------------------
# Per-point colour resolution for bar/column
# ---------------------------------------------------------------------------

def _point_color(dp, highlight_lower: str, accent: RGBColor,
                 muted: RGBColor) -> RGBColor:
    """Return colour for a single bar/column point."""
    if dp.color:
        c = _hex_to_rgb(dp.color)
        if c:
            return c
    if highlight_lower and dp.label.lower() == highlight_lower:
        return accent
    if highlight_lower:
        return muted
    return accent   # no highlight → all bars same accent colour


# ---------------------------------------------------------------------------
# Native chart renderers
# ---------------------------------------------------------------------------

def render_bar_chart(slide, left, top, width, height, spec: ChartSpec):
    """Horizontal clustered bar chart — PBI-style per-point colouring."""
    if not spec.data:
        return None

    cd = ChartData()
    cd.categories = [dp.label for dp in spec.data]
    cd.add_series("", [dp.value for dp in spec.data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    series = co.series[0]
    hl = (spec.highlight or '').lower()
    for idx, dp in enumerate(spec.data):
        try:
            pt = series.points[idx]
            c  = _point_color(dp, hl, PBI_BLUE, PBI_BLUE_LIGHT)
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = c
        except Exception:
            pass

    _set_data_labels(series, font_size_pt=9, color=DARK_GRAY,
                     show_value=True, position='outEnd')
    return cf


def render_bar_stacked_chart(slide, left, top, width, height, spec: ChartSpec):
    if not spec.series:
        return render_bar_chart(slide, left, top, width, height, spec)

    cd  = ChartData()
    cats = list({s.get('label','') for s in spec.series})
    cats = sorted(cats, key=lambda x: next((i for i,s in enumerate(spec.series) if s.get('label')==x), 0))
    cd.categories = cats

    series_map: dict = {}
    for item in spec.series:
        sn = item.get('series','')
        if sn not in series_map:
            series_map[sn] = {c: 0.0 for c in cats}
        series_map[sn][item.get('label','')] = float(item.get('value',0))

    for sn, cv in series_map.items():
        cd.add_series(sn, [cv.get(c,0.0) for c in cats])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=False)

    for i, s in enumerate(co.series):
        _set_series_color(s, MULTI_PALETTE[i % len(MULTI_PALETTE)])
    _add_legend(co)
    return cf


def render_column_chart(slide, left, top, width, height, spec: ChartSpec):
    """Vertical clustered column chart."""
    if not spec.data:
        return None

    cd = ChartData()
    cd.categories = [dp.label for dp in spec.data]
    cd.add_series("", [dp.value for dp in spec.data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    series = co.series[0]
    hl = (spec.highlight or '').lower()
    for idx, dp in enumerate(spec.data):
        try:
            pt = series.points[idx]
            c  = _point_color(dp, hl, PBI_BLUE, PBI_BLUE_LIGHT)
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = c
        except Exception:
            pass

    _set_data_labels(series, font_size_pt=8, color=DARK_GRAY,
                     show_value=True, position='outEnd')
    return cf


def render_column_stacked_chart(slide, left, top, width, height, spec: ChartSpec):
    if not spec.series:
        return render_column_chart(slide, left, top, width, height, spec)

    cd   = ChartData()
    cats = list({s.get('label','') for s in spec.series})
    cats = sorted(cats, key=lambda x: next((i for i,s in enumerate(spec.series) if s.get('label')==x), 0))
    cd.categories = cats

    series_map: dict = {}
    for item in spec.series:
        sn = item.get('series','')
        if sn not in series_map:
            series_map[sn] = {c: 0.0 for c in cats}
        series_map[sn][item.get('label','')] = float(item.get('value',0))

    for sn, cv in series_map.items():
        cd.add_series(sn, [cv.get(c,0.0) for c in cats])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, s in enumerate(co.series):
        _set_series_color(s, MULTI_PALETTE[i % len(MULTI_PALETTE)])
    _add_legend(co)
    return cf


def render_line_chart(slide, left, top, width, height, spec: ChartSpec):
    """Line chart with markers — multi-series, PBI-styled lines."""
    if not spec.series:
        return None

    cd   = ChartData()
    pts0 = spec.series[0].get('points', [])
    cd.categories = [str(p.get('x','')) for p in pts0]

    for s in spec.series:
        cd.add_series(s.get('name',''), [float(p.get('y',0)) for p in s.get('points',[])])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, series in enumerate(co.series):
        s_spec = spec.series[i] if i < len(spec.series) else {}
        # Colour from spec or palette
        c_str  = s_spec.get('color','')
        c      = _hex_to_rgb(c_str) if c_str else MULTI_PALETTE[i % len(MULTI_PALETTE)]
        try:
            series.format.line.color.rgb = c
            series.format.line.width     = Pt(2.0)
            series.marker.style          = 4      # circle
            series.marker.size           = 6
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = c
            series.marker.format.line.fill.background()
        except Exception:
            pass

    if len(spec.series) > 1:
        _add_legend(co, XL_LEGEND_POSITION.BOTTOM, font_size=9)

    return cf


def render_area_chart(slide, left, top, width, height, spec: ChartSpec):
    if not spec.series:
        return None

    cd   = ChartData()
    pts0 = spec.series[0].get('points', [])
    cd.categories = [str(p.get('x','')) for p in pts0]

    for s in spec.series:
        cd.add_series(s.get('name',''), [float(p.get('y',0)) for p in s.get('points',[])])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.AREA, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, series in enumerate(co.series):
        _set_series_color(series, MULTI_PALETTE[i % len(MULTI_PALETTE)])

    if len(spec.series) > 1:
        _add_legend(co, XL_LEGEND_POSITION.BOTTOM, font_size=9)

    return cf


def render_pie_chart(slide, left, top, width, height, spec: ChartSpec):
    if not spec.data:
        return None

    cd = ChartData()
    cd.categories = [dp.label for dp in spec.data]
    cd.add_series("", [dp.value for dp in spec.data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=False)

    series = co.series[0]
    for idx, dp in enumerate(spec.data):
        try:
            pt = series.points[idx]
            c  = _hex_to_rgb(dp.color) if dp.color else MULTI_PALETTE[idx % len(MULTI_PALETTE)]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = c
        except Exception:
            pass

    _set_data_labels(series, 9, DARK_GRAY)
    _add_legend(co, XL_LEGEND_POSITION.RIGHT, font_size=9)
    return cf


def render_donut_chart(slide, left, top, width, height, spec: ChartSpec):
    """Doughnut chart — habit-tier palette, legend on right."""
    if not spec.data:
        return None

    cd = ChartData()
    cd.categories = [dp.label for dp in spec.data]
    cd.add_series("", [dp.value for dp in spec.data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=False)

    # Default palette for 4-slice habit charts: gray→light→mid→dark
    default_palette = HABIT_M365  # overridden by dp.color when provided

    series = co.series[0]
    for idx, dp in enumerate(spec.data):
        try:
            pt = series.points[idx]
            c  = _hex_to_rgb(dp.color) if dp.color else default_palette[idx % len(default_palette)]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = c
        except Exception:
            pass

    _set_data_labels(series, 9, WHITE, show_value=False,
                     show_percent=True, position='ctr')
    _add_legend(co, XL_LEGEND_POSITION.RIGHT, font_size=9)
    return cf


def render_scatter_chart(slide, left, top, width, height, spec: ChartSpec):
    """
    XY scatter — each data point is its own series for per-point labels/colors.
    Supports: color per point, highlight, x_label / y_label axis titles.
    Data labels show the series name (org / entity name).
    """
    if not spec.series:
        return None

    cd = XyChartData()
    for s in spec.series:
        sd = cd.add_series(s.get('name', ''))
        sd.add_data_point(float(s.get('x', 0)), float(s.get('y', 0)))

    cf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, series in enumerate(co.series):
        s_spec = spec.series[i] if i < len(spec.series) else {}
        is_hl  = s_spec.get('highlight', False)
        c_str  = s_spec.get('color', '')
        c      = _hex_to_rgb(c_str) if c_str else (PBI_PINK if is_hl else PBI_PINK_LIGHT)
        sz     = 12 if is_hl else 8
        try:
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = c
            series.marker.format.line.fill.background()
            series.marker.size = sz
        except Exception:
            pass
        # Data label: show org name above the marker
        _set_data_labels(series, font_size_pt=8, color=DARK_GRAY,
                         show_value=False, show_series_name=True, position='t')

    # Axis titles
    if spec.x_label:
        try:
            _set_axis_title(co.value_axis, spec.x_label, font_size_pt=8)
        except Exception:
            pass
    if spec.y_label:
        try:
            _set_axis_title(co.category_axis, spec.y_label, font_size_pt=8)
        except Exception:
            pass

    # Show axis tick labels at 8pt
    for ax in [co.value_axis, co.category_axis]:
        try:
            ax.tick_labels.font.size      = Pt(8)
            ax.tick_labels.font.name      = FONT_NAME
            ax.tick_labels.font.color.rgb = MID_GRAY
        except Exception:
            pass

    # No legend — labels on the points are sufficient
    try:
        co.has_legend = False
    except Exception:
        pass

    return cf


def render_bubble_chart(slide, left, top, width, height, spec: ChartSpec):
    """
    Bubble chart — each bubble is its own series for per-bubble color + labels.
    spec.series entries: {name, x, y, size (bubble area), color (hex), highlight}
    Axis titles from spec.x_label / spec.y_label.
    Data labels show the entity name.
    """
    if not spec.series:
        return None

    cd = BubbleChartData()
    for s in spec.series:
        sd = cd.add_series(s.get('name', ''))
        # Normalise bubble size to a readable range
        raw_size = float(s.get('size', 5))
        sd.add_data_point(float(s.get('x', 0)), float(s.get('y', 0)), raw_size)

    cf = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, series in enumerate(co.series):
        s_spec = spec.series[i] if i < len(spec.series) else {}
        is_hl  = s_spec.get('highlight', False)
        c_str  = s_spec.get('color', '')
        c      = _hex_to_rgb(c_str) if c_str else (PBI_PINK if is_hl else PBI_PINK_LIGHT)
        _set_series_color(series, c)
        # Semi-transparent border that matches bubble color for definition
        try:
            series.format.line.color.rgb = c
            series.format.line.width     = Pt(0.75)
        except Exception:
            pass
        # Data label: entity name above bubble
        _set_data_labels(series, font_size_pt=8, color=DARK_GRAY,
                         show_value=False, show_series_name=True, position='t')

    # Axis titles
    if spec.x_label:
        try:
            _set_axis_title(co.value_axis, spec.x_label, font_size_pt=8)
        except Exception:
            pass
    if spec.y_label:
        try:
            _set_axis_title(co.category_axis, spec.y_label, font_size_pt=8)
        except Exception:
            pass

    # Axis tick labels visible
    for ax in [co.value_axis, co.category_axis]:
        try:
            ax.tick_labels.font.size      = Pt(8)
            ax.tick_labels.font.name      = FONT_NAME
            ax.tick_labels.font.color.rgb = MID_GRAY
        except Exception:
            pass

    try:
        co.has_legend = False
    except Exception:
        pass

    return cf


def render_radar_chart(slide, left, top, width, height, spec: ChartSpec):
    if not spec.series:
        return None

    cd   = ChartData()
    pts0 = spec.series[0].get('points', [])
    cd.categories = [str(p.get('x','')) for p in pts0]

    for s in spec.series:
        cd.add_series(s.get('name',''), [float(p.get('y',0)) for p in s.get('points',[])])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.RADAR_FILLED, left, top, width, height, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=True)

    for i, series in enumerate(co.series):
        _set_series_color(series, MULTI_PALETTE[i % len(MULTI_PALETTE)])

    if len(spec.series) > 1:
        _add_legend(co, XL_LEGEND_POSITION.BOTTOM, font_size=9)

    return cf


# ---------------------------------------------------------------------------
# Shape-based renderers
# ---------------------------------------------------------------------------

def render_kpi_card(slide, left, top, width, height, spec: ChartSpec):
    """Large single-number KPI card — clean, borderless."""
    # Subtle background
    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 253)
    bg.line.fill.background()

    val_h   = height * 0.50
    lbl_h   = height * 0.28
    sub_h   = height * 0.18

    # Auto-scale value font to fill available space
    h_pt    = height / 12700           # EMU → pt (approx)
    val_pts = max(20, min(44, int(h_pt * 0.40)))

    # Value
    _add_textbox(slide, spec.value,
                 left + Inches(0.08), top + height * 0.06,
                 width - Inches(0.16), val_h,
                 font_size_pt=val_pts, bold=True,
                 color=PBI_BLUE, align=PP_ALIGN.CENTER)

    # Label
    lbl_pts = max(9, min(14, int(h_pt * 0.11)))
    _add_textbox(slide, spec.label,
                 left + Inches(0.08), top + height * 0.06 + val_h,
                 width - Inches(0.16), lbl_h,
                 font_size_pt=lbl_pts, bold=False,
                 color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Subtitle
    if spec.subtitle:
        sub_pts = max(8, min(11, int(h_pt * 0.09)))
        _add_textbox(slide, spec.subtitle,
                     left + Inches(0.08), top + height * 0.06 + val_h + lbl_h,
                     width - Inches(0.16), sub_h,
                     font_size_pt=sub_pts, bold=False,
                     color=MID_GRAY, align=PP_ALIGN.CENTER)
    return bg


def render_heatmap(slide, left, top, width, height, spec: ChartSpec):
    rows = spec.rows or []
    cols = spec.columns or []
    vals = spec.values or []

    if not rows or not cols or not vals:
        return None

    n_rows, n_cols = len(rows), len(cols)
    row_lbl_w  = min(Inches(1.0), width * 0.22)
    header_h   = min(Inches(0.28), height * 0.12)
    cell_w     = (width  - row_lbl_w) / n_cols
    cell_h     = (height - header_h)  / n_rows
    all_vals   = [v for row in vals for v in row if v is not None]
    max_val    = max(all_vals) if all_vals else 1.0

    # Column headers
    for ci, col_name in enumerate(cols):
        _add_textbox(slide, col_name,
                     left + row_lbl_w + ci * cell_w, top,
                     cell_w, header_h,
                     font_size_pt=8, bold=True, color=DARK_GRAY,
                     align=PP_ALIGN.CENTER)

    for ri, row_name in enumerate(rows):
        ry = top + header_h + ri * cell_h
        _add_textbox(slide, row_name,
                     left, ry, row_lbl_w, cell_h,
                     font_size_pt=8, bold=False, color=DARK_GRAY,
                     align=PP_ALIGN.RIGHT)
        row_vs = vals[ri] if ri < len(vals) else []
        for ci in range(n_cols):
            cx   = left + row_lbl_w + ci * cell_w
            val  = row_vs[ci] if ci < len(row_vs) else 0.0
            norm = val / max_val if max_val > 0 else 0.0
            r    = int(255 + (0   - 255) * norm)
            g    = int(255 + (114 - 255) * norm)
            b    = int(255 + (198 - 255) * norm)
            cell_c = RGBColor(max(0,min(255,r)), max(0,min(255,g)), max(0,min(255,b)))

            cell = slide.shapes.add_shape(1, cx, ry, cell_w, cell_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = cell_c
            cell.line.color.rgb      = WHITE
            cell.line.width          = Pt(0.5)

            cell_h_pt  = cell_h / 12700
            fs         = max(6, min(10, int(cell_h_pt * 0.35)))
            txt_color  = WHITE if norm > 0.55 else DARK_GRAY
            fmt        = spec.value_format or '.1f'
            try:
                lbl = format(val, fmt)
            except Exception:
                lbl = str(val)
            _add_textbox(slide, lbl,
                         cx + Inches(0.02), ry + Inches(0.02),
                         cell_w - Inches(0.04), cell_h - Inches(0.04),
                         font_size_pt=fs, color=txt_color,
                         align=PP_ALIGN.CENTER)
    return None


def render_table(slide, left, top, width, height, spec: ChartSpec):
    cols = spec.table_columns or []
    rows = spec.table_rows    or []
    if not cols or not rows:
        return None

    n_c, n_r = len(cols), len(rows)
    tbl      = slide.shapes.add_table(n_r + 1, n_c, left, top, width, height).table

    col_w    = width  // n_c
    row_h    = height // (n_r + 1)
    for i in range(n_c):
        tbl.columns[i].width = col_w
    for i in range(n_r + 1):
        tbl.rows[i].height = row_h

    row_h_pt = row_h / 12700
    fs       = max(7, min(11, int(row_h_pt * 0.38)))

    # Header
    for ci, col_name in enumerate(cols):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = PBI_BLUE_DARK
        tf  = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs
               else tf.paragraphs[0].add_run())
        run.font.name  = FONT_NAME
        run.font.size  = Pt(fs)
        run.font.bold  = True
        run.font.color.rgb = WHITE
        for attr in ('margin_left','margin_right','margin_top','margin_bottom'):
            try: setattr(cell, attr, 0)
            except: pass

    ALT = RGBColor(242, 246, 252)
    for ri, row in enumerate(rows):
        bg = ALT if ri % 2 == 0 else WHITE
        for ci in range(n_c):
            cell = tbl.cell(ri + 1, ci)
            val  = row[ci] if ci < len(row) else ''
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf  = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            is_bold = (ci == spec.highlight_col)
            run = (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs
                   else tf.paragraphs[0].add_run())
            run.font.name  = FONT_NAME
            run.font.size  = Pt(fs)
            run.font.bold  = is_bold
            run.font.color.rgb = PBI_BLUE_DARK if is_bold else DARK_GRAY
            for attr in ('margin_left','margin_right','margin_top','margin_bottom'):
                try: setattr(cell, attr, 0)
                except: pass
    return tbl


def render_funnel(slide, left, top, width, height, spec: ChartSpec):
    data = spec.data or []
    if not data:
        return None

    n         = len(data)
    stage_h   = height / n
    max_val   = max(dp.value for dp in data) if data else 1.0
    lbl_w     = width * 0.28
    val_w     = width * 0.18
    bar_area  = width - lbl_w - val_w
    bar_off   = lbl_w

    # Platform-aware palette
    palette = [PBI_BLUE_DARK, PBI_BLUE, PBI_BLUE_MID, PBI_BLUE_LIGHT, LIGHT_GRAY]

    for i, dp in enumerate(data):
        frac  = dp.value / max_val if max_val > 0 else 0
        bar_w = bar_area * frac
        bar_x = left + bar_off + (bar_area - bar_w) / 2
        bar_y = top + i * stage_h
        pad   = Inches(0.06)

        c = _hex_to_rgb(dp.color) if dp.color else palette[i % len(palette)]
        bar = slide.shapes.add_shape(1, bar_x, bar_y + pad, bar_w, stage_h - 2*pad)
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

        _add_textbox(slide, dp.label,
                     left, bar_y + pad, lbl_w - Inches(0.08), stage_h - 2*pad,
                     font_size_pt=9, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

        try:
            vt = format(dp.value, '.0f') if dp.value == int(dp.value) else format(dp.value, '.1f')
        except Exception:
            vt = str(dp.value)
        _add_textbox(slide, vt,
                     left + lbl_w + bar_area + Inches(0.05), bar_y + pad,
                     val_w, stage_h - 2*pad,
                     font_size_pt=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.LEFT)
    return None


def _squarify(items, x, y, w, h):
    """Simple squarified treemap layout → list of (label, val, x, y, w, h)."""
    if not items:
        return []
    total = sum(v for _, v in items)
    if total <= 0 or w <= 0 or h <= 0:
        return []
    result = []

    def layout_row(row, x0, y0, rw, rh, horiz):
        rs = sum(v for _,v in row)
        if rs <= 0: return
        off = 0
        for lbl, val in row:
            f = val / rs
            if horiz:
                result.append((lbl, val, x0, y0+off, rw, rh*f)); off += rh*f
            else:
                result.append((lbl, val, x0+off, y0, rw*f, rh)); off += rw*f

    def recurse(items, x0, y0, w0, h0):
        if not items: return
        tot = sum(v for _,v in items)
        if tot <= 0: return
        horiz = w0 <= h0
        cur, best = [], float('inf')
        for item in items:
            trial = cur + [item]
            rs    = sum(v for _,v in trial)
            frac  = rs / tot
            if horiz:
                rh = h0 * frac
                ratios = [(v/rs*rh/w0 if w0 > 0 else 1e9) for _,v in trial]
            else:
                rw = w0 * frac
                ratios = [(v/rs*rw/h0 if h0 > 0 else 1e9) for _,v in trial]
            mr = max((max(r,1/r) if r > 0 else 1e9) for r in ratios)
            if mr <= best:
                best, cur = mr, trial
            else:
                break
        if not cur:
            cur = [items[0]]
        rs   = sum(v for _,v in cur)
        frac = rs / tot
        if horiz:
            rh = h0 * frac
            layout_row(cur, x0, y0, w0, rh, True)
            recurse([i for i in items if i not in cur], x0, y0+rh, w0, h0-rh)
        else:
            rw = w0 * frac
            layout_row(cur, x0, y0, rw, h0, False)
            recurse([i for i in items if i not in cur], x0+rw, y0, w0-rw, h0)

    recurse(sorted(items, key=lambda x: x[1], reverse=True), x, y, w, h)
    return result


def render_treemap(slide, left, top, width, height, spec: ChartSpec):
    data = spec.data or []
    if not data:
        return None

    items  = [(dp.label, dp.value) for dp in data]
    rects  = _squarify(items, left, top, width, height)
    palette= [PBI_BLUE_DARK, PBI_BLUE, PBI_BLUE_MID, PBI_PURPLE, PBI_PINK_DARK,
              PBI_PURPLE_MID, PBI_BLUE_LIGHT]
    dmap   = {dp.label: dp for dp in data}

    for i, (lbl, val, rx, ry, rw, rh) in enumerate(rects):
        dp = dmap.get(lbl)
        c  = _hex_to_rgb(dp.color) if (dp and dp.color) else palette[i % len(palette)]
        gap = Inches(0.025)
        sh  = slide.shapes.add_shape(
            1,
            rx + gap, ry + gap,
            max(rw - 2*gap, Inches(0.05)),
            max(rh - 2*gap, Inches(0.05))
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()

        rw_pt = rw / 12700
        rh_pt = rh / 12700
        fs    = max(6, min(11, int(min(rw_pt, rh_pt) * 0.22)))
        if fs >= 6:
            _add_textbox(slide, lbl,
                         rx + gap + Inches(0.05), ry + gap + Inches(0.05),
                         max(rw - 2*gap - Inches(0.1), Inches(0.05)),
                         max(rh - 2*gap - Inches(0.1), Inches(0.05)),
                         font_size_pt=fs, color=WHITE, align=PP_ALIGN.LEFT)
    return None


def render_gauge(slide, left, top, width, height, spec: ChartSpec):
    try:
        raw = float(spec.value or 0)
    except (ValueError, TypeError):
        raw = 0.0
    max_val = spec.max_value or 100.0
    frac    = max(0.0, min(1.0, raw / max_val))
    filled  = frac  * 50.0
    empty   = (1-frac) * 50.0
    hidden  = 50.0

    cd = ChartData()
    cd.categories = ['Filled', 'Empty', 'Hidden']
    cd.add_series('', [filled, empty, hidden])

    chart_h = height * 1.35
    chart_t = top - height * 0.18

    cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, chart_t, width, chart_h, cd)
    co = cf.chart
    _apply_pbi_chart_style(co, cf, show_gridlines=False)
    try:
        co.has_legend = False
    except Exception:
        pass

    series = co.series[0]
    for idx, c in enumerate([PBI_BLUE, GRID_GRAY, WHITE]):
        try:
            pt = series.points[idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = c
            if idx == 2:
                pt.format.line.fill.background()
        except Exception:
            pass

    val_txt = f"{int(round(raw))}{'%' if max_val == 100 else ''}"
    _add_textbox(slide, val_txt,
                 left, top + height * 0.44, width, height * 0.28,
                 font_size_pt=18, bold=True, color=PBI_BLUE, align=PP_ALIGN.CENTER)
    if spec.label:
        _add_textbox(slide, spec.label,
                     left, top + height * 0.68, width, height * 0.22,
                     font_size_pt=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    return cf


def render_kpi_row(slide, left, top, width, height, spec: ChartSpec):
    """
    A horizontal row of 2-4 compact KPI tiles — suitable when multiple headline
    numbers need to be shown together.  Pass KPI list via spec.series:
      [{"value": "79%", "label": "Return Rate", "subtitle": "Target >50%"}, ...]
    """
    kpis = spec.series or []
    if not kpis:
        return render_kpi_card(slide, left, top, width, height, spec)

    n      = len(kpis)
    gap    = Inches(0.10)
    tile_w = int((width - gap * (n - 1)) // n)
    h_pt   = height / 12700
    val_pts = max(22, min(40, int(h_pt * 0.38)))
    lbl_pts = max(9,  min(13, int(h_pt * 0.11)))
    sub_pts = max(8,  min(10, int(h_pt * 0.09)))

    for i, kpi in enumerate(kpis):
        tx = int(left + i * (tile_w + gap))

        # Tile background with light border
        bg = slide.shapes.add_shape(1, tx, top, tile_w, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(248, 250, 253)
        bg.line.color.rgb      = RGBColor(210, 210, 210)
        bg.line.width          = Pt(0.5)

        # Value
        _add_textbox(slide, kpi.get('value', ''),
                     tx + Inches(0.08), top + int(height * 0.06),
                     tile_w - Inches(0.16), int(height * 0.48),
                     font_size_pt=val_pts, bold=True,
                     color=PBI_BLUE, align=PP_ALIGN.CENTER)

        # Label
        _add_textbox(slide, kpi.get('label', ''),
                     tx + Inches(0.08), top + int(height * 0.06 + height * 0.48),
                     tile_w - Inches(0.16), int(height * 0.30),
                     font_size_pt=lbl_pts, bold=False,
                     color=DARK_GRAY, align=PP_ALIGN.CENTER)

        # Subtitle
        if kpi.get('subtitle'):
            _add_textbox(slide, kpi['subtitle'],
                         tx + Inches(0.08),
                         top + int(height * 0.06 + height * 0.48 + height * 0.30),
                         tile_w - Inches(0.16), int(height * 0.20),
                         font_size_pt=sub_pts, bold=False,
                         color=MID_GRAY, align=PP_ALIGN.CENTER)

    return None


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def render_chart(slide, left, top, width, height, spec: ChartSpec):
    """Dispatch to the correct renderer. Returns shape/chart_frame or None."""
    dispatch = {
        'bar':              render_bar_chart,
        'bar_stacked':      render_bar_stacked_chart,
        'column':           render_column_chart,
        'column_stacked':   render_column_stacked_chart,
        'line':             render_line_chart,
        'area':             render_area_chart,
        'pie':              render_pie_chart,
        'donut':            render_donut_chart,
        'scatter':          render_scatter_chart,
        'bubble':           render_bubble_chart,
        'radar':            render_radar_chart,
        'kpi':              render_kpi_card,
        'kpi_row':          render_kpi_row,
        'heatmap':          render_heatmap,
        'table':            render_table,
        'funnel':           render_funnel,
        'treemap':          render_treemap,
        'gauge':            render_gauge,
    }
    fn = dispatch.get((spec.type or '').lower().strip())
    if fn is None:
        return None
    try:
        return fn(slide, left, top, width, height, spec)
    except Exception as e:
        print(f"  WARNING: chart render failed ({spec.type}): {e}")
        return None
