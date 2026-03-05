#!/usr/bin/env python3
"""
Dashboard-to-Executive-Deck Converter

Works with Claude Code or GitHub Copilot Chat to generate analyst-grade insights.
No API key needed — leverages the interactive assistant session.

Workflow:
1. Extract slides as images (deterministic)
2. Prepare analysis request for the assistant
3. Assistant analyzes each image and generates insights
4. Build final presentation with insights
"""

import argparse
import json
import sys
import time
import os
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io
from lib.extraction.pdf_extractor import prepare_pdf_for_analysis


def extract_slide_as_image(prs, slide_idx, output_path):
    """Extract a slide's dashboard image"""
    slide = prs.slides[slide_idx]

    # Find the main dashboard image
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            image_stream = io.BytesIO(shape.image.blob)
            img = Image.open(image_stream)
            img.save(output_path)
            return True

    return False


def extract_slide_title(slide):
    """Extract slide title, removing emojis"""
    import re

    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            title = shape.text.strip()
            # Remove emoji characters
            title = re.sub(r'[\U00010000-\U0010ffff]', '', title).strip()
            return title

    return "Untitled Slide"


def classify_slide_type(title):
    """Classify slide type for context"""
    title_lower = title.lower()

    if 'trend' in title_lower or 'over time' in title_lower:
        return 'trends'
    elif 'leaderboard' in title_lower or 'top' in title_lower:
        return 'leaderboard'
    elif 'health' in title_lower or 'overview' in title_lower:
        return 'health_check'
    elif 'habit' in title_lower or 'frequency' in title_lower:
        return 'habit_formation'
    elif 'license' in title_lower or 'priority' in title_lower:
        return 'license_priority'
    else:
        return 'general'


def detect_file_type(file_path):
    """
    Detect if file is PPTX, PDF, or PBIP based on extension / directory contents.

    Args:
        file_path: Path to input file or directory

    Returns:
        'pptx', 'pdf', or 'pbip'

    Raises:
        ValueError: If file type is unsupported
    """
    p = Path(file_path)
    suffix = p.suffix.lower()

    if suffix == '.pptx':
        return 'pptx'
    elif suffix == '.pdf':
        return 'pdf'
    elif suffix == '.pbip':
        return 'pbip'
    elif suffix == '.pbix':
        return 'pbix'
    elif p.is_dir() and any(p.glob('*.pbip')):
        return 'pbip'
    else:
        raise ValueError(
            f"Unsupported file type: {suffix}. Supported formats: .pptx, .pdf, .pbip, .pbix"
        )


def _check_pbi_mcp_setup() -> bool:
    """
    Warn clearly if the Power BI Modeling MCP is not installed/configured.
    Returns True if MCP is ready for DAX queries, False otherwise.
    Does NOT block execution.
    """
    mcp_json = Path(".mcp.json")

    configured = False
    exe_valid  = False
    if mcp_json.exists():
        try:
            cfg    = json.loads(mcp_json.read_text(encoding="utf-8"))
            server = cfg.get("mcpServers", {}).get("powerbi-modeling")
            if server:
                configured = True
                exe_valid  = Path(server.get("command", "")).exists()
        except Exception:
            pass

    if configured and exe_valid:
        print("  OK Power BI Modeling MCP configured — DAX query mode enabled")
        return True

    print()
    print("  " + "!" * 66)
    if not configured:
        print("  !! Power BI Modeling MCP is NOT installed")
        print("  !!")
        print("  !! No live DAX queries will run. Claude will use screenshots only.")
        print("  !! To enable exact DAX values, install the MCP first:")
        print("  !!")
        print("  !!     python setup_pbi_mcp.py")
        print("  !!")
        print("  !! Then restart Claude Code and re-run this command.")
    else:
        print("  !! Power BI MCP is configured in .mcp.json but the executable")
        print("  !! was not found at the registered path.")
        print("  !! Re-run the setup script to fix the path:")
        print("  !!")
        print("  !!     python setup_pbi_mcp.py --force")
    print("  " + "!" * 66)
    print()
    print("  MCP not available — falling back to image-only analysis.")
    print()
    return False


def _is_mcp_ready() -> bool:
    """Silent MCP check — returns True/False without printing anything."""
    mcp_json = Path(".mcp.json")
    if not mcp_json.exists():
        return False
    try:
        cfg    = json.loads(mcp_json.read_text(encoding="utf-8"))
        server = cfg.get("mcpServers", {}).get("powerbi-modeling")
        if server and Path(server.get("command", "")).exists():
            return True
    except Exception:
        pass
    return False


def prepare_for_analysis(source_path, use_text_layer: bool = False):
    """
    Extract slides/pages and prepare analysis request for the assistant.
    Supports both .pptx and .pdf input files.

    Args:
        source_path: Path to source file (.pptx or .pdf)

    Returns:
        Path to analysis_request.json file
    """
    file_type = detect_file_type(source_path)

    if file_type == 'pdf':
        return prepare_pdf_for_analysis(source_path, use_text_layer=use_text_layer)

    if file_type == 'pbip':
        _check_pbi_mcp_setup()
        from lib.extraction.pbip_extractor import prepare_pbip_for_analysis
        return prepare_pbip_for_analysis(source_path)

    if file_type == 'pbix':
        _check_pbi_mcp_setup()
        from lib.extraction.pbix_extractor import prepare_pbix_for_analysis
        return prepare_pbix_for_analysis(source_path)

    # Original PPTX logic continues below
    print("=" * 70)
    print("PREPARING SLIDES FOR ANALYSIS")
    print("=" * 70)

    prs = Presentation(source_path)

    # Create temp directory for images
    Path('temp').mkdir(exist_ok=True)

    slides_to_analyze = []

    print(f"\nExtracting {len(prs.slides)} slides...")
    print("  Skipping slide 1 (title page)...")

    for idx, slide in enumerate(prs.slides):
        # Skip the first slide (title page with metadata, no insights)
        if idx == 0:
            continue

        title = extract_slide_title(slide)
        image_path = f"temp/slide_{idx+1}.png"

        # Extract dashboard image
        has_image = extract_slide_as_image(prs, idx, image_path)

        if has_image:
            slide_info = {
                'slide_number': idx + 1,
                'title': title,
                'image_path': image_path,
                'slide_type': classify_slide_type(title)
            }

            slides_to_analyze.append(slide_info)
            print(f"  OK Slide {idx+1}: {title[:50]}...")

    if use_text_layer:
        from lib.extraction.text_layer_extractor import enrich_slides_with_pptx_text
        print("\nExtracting text layer from PPTX...")
        enrich_slides_with_pptx_text(source_path, slides_to_analyze)

    # OCR fallback: for slides where text-layer was insufficient (e.g. PBI
    # dashboard PNGs), run EasyOCR to extract real numbers and labels.
    if use_text_layer:
        try:
            from lib.extraction.ocr_extractor import enrich_slides_with_ocr
            print("\nChecking text quality — OCR fallback if needed...")
            ocr_count = enrich_slides_with_ocr(slides_to_analyze)
            if ocr_count:
                print(f"  OK  OCR enriched {ocr_count} slide(s)")
            else:
                print("  OK  Text layer sufficient — no OCR needed")
        except ImportError:
            print("\n  WARN  easyocr not installed — skipping OCR fallback")
            print("         Install with: pip install easyocr")

    # Save analysis request
    request_file = 'temp/analysis_request.json'
    with open(request_file, 'w', encoding='utf-8') as f:
        json.dump({
            'source_file': source_path,
            'source_type': 'pptx',
            'total_slides': len(slides_to_analyze),
            'slides': slides_to_analyze,
            'text_layer_used': use_text_layer,
        }, f, indent=2)

    print(f"\nOK Prepared {len(slides_to_analyze)} slides for analysis")
    print(f"OK Analysis request saved to: {request_file}")

    return request_file


def show_claude_instructions(request_file, context=None):
    """Show instructions for Claude to generate insights (manual mode)"""

    print("\n" + "=" * 70)
    print("NEXT STEP: CLAUDE GENERATES INSIGHTS")
    print("=" * 70)

    context_line = f'\n    Focus: {context}' if context else ''
    print(f"""
Claude will now analyze each dashboard image and generate insights.

Please say to Claude Code:

    "Analyze the dashboards in {request_file} and generate analyst-grade insights.{context_line}
    Save results to temp/insights.json"

After Claude completes, run:
    python convert_dashboard.py --build --output [output.pptx]
""")


def show_copilot_instructions(request_file, context=None):
    """Show instructions for Copilot Chat to generate insights.

    When Copilot Chat (agent mode) is running this command, it will see
    these instructions in the terminal output, read the extracted data,
    generate insights, write the JSON, and then run the build step — all
    automatically within the same session.
    """
    print("\n" + "=" * 70)
    print("COPILOT CHAT: GENERATE INSIGHTS")
    print("=" * 70)

    with open(request_file, 'r', encoding='utf-8') as f:
        request = json.load(f)

    has_pbip_context = Path('temp/pbip_context.json').exists()
    is_pbip = request.get('source_type') in ('pbip', 'pbix') or has_pbip_context

    context_line = f"\n    Focus: {context}" if context else ""

    print("\n--- COPILOT: do the following now ---")
    if is_pbip:
        print("  1. Read temp/pbip_context.json and temp/analysis_request.json")
        print("  2. For each page, execute DAX queries via the powerbi-modeling MCP")
        print("  3. Generate analyst-grade insights from queried values (see COPILOT.md)")
        print("  4. Write temp/insights.json (schema in COPILOT.md)")
    else:
        print(f"  1. Read {request_file}")
        print("  2. Read each slide image listed in analysis_request.json")
        print("  3. Use text_layer + text_metrics as supporting data; verify against images")
        print("  4. Generate analyst-grade insights (see COPILOT.md for formula & schema)")
        print("  5. Write temp/insights.json")

    print(f"\n  Then build the deck:{context_line}")
    print("    python convert_dashboard.py --build --output <output>.pptx")
    print("\n--- end instructions ---")


def _resolve_assistant(selection: str) -> str:
    """Resolve assistant based on explicit selection or environment."""
    if selection in ("claude", "copilot"):
        return selection

    # Best-effort detection for Claude Code environments.
    claude_markers = [
        "CLAUDE_CODE",
        "ANTHROPIC_API_KEY",
        "CLAUDE_SESSION",
    ]
    for key in claude_markers:
        if os.environ.get(key):
            return "claude"

    return "copilot"


def trigger_claude_analysis(request_file, context=None):
    """Trigger Claude Code to analyze dashboards (automatic mode)"""

    print("\n" + "=" * 70)
    print("STEP 2: CLAUDE ANALYSIS")
    print("=" * 70)

    # Load request to show Claude what to analyze
    with open(request_file, 'r', encoding='utf-8') as f:
        request = json.load(f)

    is_pbip = request.get('source_type') in ('pbip', 'pbix') or Path('temp/pbip_context.json').exists()

    if is_pbip and _is_mcp_ready():
        _trigger_pbip_analysis(request, context=context)
    elif is_pbip and not _is_mcp_ready():
        _trigger_image_analysis(request, mcp_missing=True, context=context)
    else:
        _trigger_image_analysis(request, context=context)


def _trigger_image_analysis(request, mcp_missing=False, context=None):
    """Show image-based analysis instructions (PPTX / PDF path, or PBIX/PBIP without MCP)."""
    if mcp_missing:
        print("\n" + "!" * 70)
        print("!! NO MCP INSTALLED — RUNNING IN IMAGE-ONLY MODE")
        print("!!")
        print("!! The Power BI Modeling MCP is not installed, so no live DAX")
        print("!! queries will run. Claude will read dashboard screenshots only.")
        print("!! Numbers will be read visually — not queried from the live model.")
        print("!!")
        print("!! To enable exact DAX values: python setup_pbi_mcp.py")
        print("!" * 70 + "\n")

    if context:
        print("=" * 70)
        print("ANALYSIS FOCUS (from --context):")
        print(f"  {context}")
        print("=" * 70 + "\n")
        print("Prioritise insights, headlines, and recommendations that directly")
        print("address the focus above. Highlight relevant numbers prominently.\n")

    print(f"\nClaude Code: Please analyze these {request['total_slides']} dashboard images.\n")

    print("For each slide:")
    for slide in request['slides']:
        print(f"  - Slide {slide['slide_number']}: {slide['title']}")
        print(f"    Image: {slide['image_path']}")
        print(f"    Type: {slide['slide_type']}")

    print("\n" + "-" * 70)
    print("CLAUDE CODE TASK:")
    print("-" * 70)
    print("""
Act as senior analyst advisor to IT decision maker.

For EACH dashboard image above:
1. Read the image file
2. Extract specific numbers visible in the dashboard
3. Generate compelling headline with specific number ([Number] + [Insight])
4. Generate 2-3 concise insights (1-2 sentences each):
   - Insight 1: Specific observation with number
   - Insight 2: Pattern/opportunity identified (include platform/feature analysis)
   - Insight 3: Actionable recommendation
5. Use friendly tone (opportunities, not criticisms)
6. If insufficient data visible, mark as "Insufficient data"

IMPORTANT: Include ALL slides in your output. Use slide_number to ensure
complete coverage. You can create insight-led titles that are better than
the source titles - just make sure every slide has an entry.

Follow Claude PowerPoint Constitution Section 5A guidelines.

Save results to: temp/insights.json

Also generate a compelling deck_title and deck_subtitle:
- deck_title: the single core story of the deck, positively framed.
  Format: "[Evocative phrase]: [Core finding or story]"
  Examples: "Copilot Impact Confirmed: $14.7M in Value with 84% Active Adoption"
            "From Reach to Routine: Building AI Habits Across the Organization"
- deck_subtitle: scope descriptor — platforms covered and time period.
  Format: "[Platform 1] · [Platform 2] · [Month Year] – [Month Year]"
  Examples: "Agents · Unlicensed Chat · M365 Copilot · Mar – Jun 2025"

CHART DATA (OPTIONAL BUT RECOMMENDED):
Each insight can include a native PowerPoint chart. When you can extract
specific data values from the dashboard, add a "chart" object alongside
the "text" field. If no data is extractable, set "chart": null.

Insight format (new, with charts):
  {"text": "Bold line || Detail text", "chart": { ... chart spec ... }}

Insight format (legacy, still accepted):
  "Bold line || Detail text"

CHART TYPE DECISION GUIDE:
  - Rankings / comparisons by category  → "bar" (horizontal) or "column" (vertical)
  - Stacked breakdown by group          → "bar_stacked" or "column_stacked"
  - Trend over time (one or more lines) → "line" or "area"
  - Part-of-whole distribution          → "pie" or "donut"
  - Single KPI number with label        → "kpi"
  - Two-axis correlation (x vs y)       → "scatter"
  - Frequency/engagement tiers          → "funnel"
  - Portfolio of items by size          → "treemap"
  - Completion % of a target            → "gauge"
  - Multi-dim comparison grid           → "heatmap"
  - Tabular leaderboard data            → "table"
  - Multi-axis capability scores        → "radar"
  - No clear chart type / no data       → null

CHART SPEC EXAMPLES (include all visible data points):

bar / column:
  {"type": "bar", "title": "Sessions by Org", "highlight": "Finance",
   "data": [{"label": "Finance", "value": 5.5}, {"label": "Legal", "value": 3.26}]}

line:
  {"type": "line", "title": "Active Users Over Time",
   "series": [
     {"name": "Agents", "points": [{"x": "Mar", "y": 180}, {"x": "Apr", "y": 220}]},
     {"name": "Chat",   "points": [{"x": "Mar", "y": 2500}, {"x": "Apr", "y": 2700}]}
   ]}

kpi:
  {"type": "kpi", "value": "4,381", "label": "Active Copilot Users", "subtitle": "across 3 platforms"}

donut:
  {"type": "donut", "title": "Usage Tiers",
   "data": [{"label": "Light", "value": 144}, {"label": "Moderate", "value": 18}, {"label": "Daily", "value": 12}]}

funnel:
  {"type": "funnel", "title": "Engagement Funnel",
   "data": [{"label": "Licensed", "value": 5000}, {"label": "Active", "value": 4381}, {"label": "Daily", "value": 12}]}

heatmap:
  {"type": "heatmap", "title": "Usage by Dept × App",
   "rows": ["Finance", "Legal"], "columns": ["Chat", "M365", "Agents"],
   "values": [[5.5, 4.3, 7.6], [3.0, 4.3, 3.5]]}

table:
  {"type": "table",
   "columns": ["Department", "Sessions/wk", "Users"],
   "rows": [["Finance", "5.5", "276"], ["Legal", "3.26", "314"]],
   "highlight_col": 0}

gauge:
  {"type": "gauge", "value": "79", "max": 100, "label": "Return Rate", "threshold": 50}

scatter:
  {"type": "scatter", "title": "Engagement vs Frequency",
   "series": [
     {"name": "Finance", "x": 5.5, "y": 4.2, "highlight": true},
     {"name": "Legal", "x": 3.26, "y": 3.1}
   ]}

treemap:
  {"type": "treemap", "title": "Agent Templates by Usage",
   "data": [{"label": "Visual Creator", "value": 297}, {"label": "Data Analyst", "value": 180}]}

ACCURACY RULES FOR CHART DATA:
- Only include data points you can point to on the dashboard
- Match units exactly (13 vs 13K vs 13M)
- For bar/column charts with >12 bars, trim to the most meaningful values
- For heatmaps with >6x6 cells, trim to most significant rows/columns
- Verify every value against its label before including

Output format:
{
  "deck_title": "Compelling story-driven title",
  "deck_subtitle": "Agents · Chat · M365 Copilot · Month Year – Month Year",
  "executive_summary": [
    "Finding with specific number -> business implication",
    ...
  ],
  "recommendations": [
    "Action: Specific recommendation with expected outcome",
    ...
  ],
  "slides": [
    {
      "slide_number": 1,
      "title": "Insight-led title (can differ from source)",
      "headline": "[Number] + [Insight]",
      "insights": [
        {
          "text": "Bold punchy line || Supporting evidence with specific data",
          "chart": {"type": "bar", "title": "...", "data": [...]}
        },
        {
          "text": "KPI insight || Detail",
          "chart": {"type": "kpi", "value": "4,381", "label": "Active Users"}
        },
        {
          "text": "Trend insight || Detail",
          "chart": null
        }
      ],
      "numbers_used": ["123", "45%"]
    }
  ]
}
""")
    print("-" * 70)


def _trigger_pbip_analysis(request, context=None):
    """Show PBIP / MCP-based analysis instructions."""
    total = request['total_slides']
    print(f"\nClaude Code: Please analyze this Power BI report ({total} pages) "
          f"using the live model via MCP.\n")

    print("Pages to analyze:")
    for slide in request['slides']:
        print(f"  - Page {slide['slide_number']}: {slide['title']} ({slide['slide_type']})")

    print("\n" + "-" * 70)
    print("CLAUDE CODE TASK (PBIP / MCP MODE):")
    print("-" * 70)

    context_steps = ""
    if context:
        context_steps = f"""
ANALYSIS FOCUS: {context}

Before running any queries, scan pbip_context.json → model.tables and
model.measures to identify which table and column best matches the focus
above (e.g. a Group, Org, Department, or Manager column).

For every pre-built query, run TWO variants:
  1. Unfiltered (baseline — as written in dax_queries)
  2. Filtered to the focus:
       CALCULATE(<measure>, '<Table>'[<Column>] = "<value>")
     or for a dimension breakdown:
       CALCULATETABLE(
           SUMMARIZECOLUMNS('<Table>'[<Column>], "<Measure>", [<Measure>]),
           '<Table>'[<FilterCol>] = "<value>"
       )

Surface the filtered results prominently in insights and headlines.
Where meaningful, state the delta vs the unfiltered baseline:
  e.g. "Group A: 48.9 actions/user vs 30.8 org average (+59%)"
"""

    print(
        "\nAct as senior analyst advisor to IT decision maker.\n\n"
        "This is a PBIP project — you have access to the LIVE Power BI model via\n"
        "the powerbi-modeling MCP server.  DO NOT estimate numbers from images.\n"
        "Query the model directly for exact values.\n"
        + (context_steps or "")
        + "\nSTEP 1: Read the context file"
    )
    print("""    Read temp/pbip_context.json
    This contains: page structure, model metadata, and pre-built DAX queries.

STEP 2: For each page, execute its DAX queries using the powerbi-modeling MCP
    - Use the execute_query (or equivalent) MCP tool
    - Power BI Desktop must be open with this project for the MCP to connect
    - Run each query in pbip_context.json → dax_queries[n].queries[m].dax
    - The returned table rows ARE your data source — use exact values

STEP 3: To drill deeper or apply filters, modify the DAX:
    - Wrap with CALCULATETABLE(..., FILTER(...)) for segment analysis
    - Use DATESYTD() / DATESINPERIOD() for time-filtered views
    - Compare filtered vs unfiltered for delta / trend insights

STEP 4: Use measure_dax fields to understand HOW each KPI is calculated:
    - If it uses DATESYTD() → it's a year-to-date figure
    - If it uses DIVIDE() → watch for zero-division handling
    - Mention the calculation context in your insights where relevant

STEP 5: Generate analyst-grade insights from QUERIED values (not estimates):
    - Every number in insights must come from an executed DAX query result
    - Include the measure name alongside the value for traceability
    - Follow the same insight formula: headline + 3 insights per slide

STEP 6: Save insights to temp/insights.json (same format as always)

IMPORTANT: Include ALL pages in your output. Use slide_number matching the
analysis_request.json to ensure complete coverage.

Also generate a compelling deck_title and deck_subtitle:
- deck_title: the single core story of the deck, positively framed.
  Format: "[Evocative phrase]: [Core finding or story]"
  Examples: "Copilot Impact Confirmed: $14.7M in Value with 84% Active Adoption"
            "From Reach to Routine: Building AI Habits Across the Organization"
- deck_subtitle: scope descriptor — platforms and time period (infer from model data if possible).
  Format: "[Platform 1] · [Platform 2] · [Month Year] – [Month Year]"
  Examples: "Agents · Unlicensed Chat · M365 Copilot · Mar – Jun 2025"

CHART DATA (OPTIONAL BUT RECOMMENDED):
Each insight can include a native PowerPoint chart. When DAX queries return
usable data, add a "chart" object alongside the "text" field.
All chart values must come from executed DAX query results only.
If no DAX data is available for a chart, set "chart": null.

CHART TYPE DECISION GUIDE:
  - Rankings / comparisons by category  → "bar" (horizontal) or "column" (vertical)
  - Stacked breakdown by group          → "bar_stacked" or "column_stacked"
  - Trend over time (one or more lines) → "line" or "area"
  - Part-of-whole distribution          → "pie" or "donut"
  - Single KPI number with label        → "kpi"
  - Two-axis correlation (x vs y)       → "scatter"
  - Frequency/engagement tiers          → "funnel"
  - Portfolio of items by size          → "treemap"
  - Completion % of a target            → "gauge"
  - Multi-dim comparison grid           → "heatmap"
  - Tabular leaderboard data            → "table"
  - No clear chart type / no data       → null

CHART SPEC EXAMPLES (all values from DAX results):

bar / column:
  {"type": "bar", "title": "Sessions by Org", "highlight": "Finance",
   "data": [{"label": "Finance", "value": 5.5}, {"label": "Legal", "value": 3.26}]}

line:
  {"type": "line", "title": "Active Users Over Time",
   "series": [
     {"name": "Agents", "points": [{"x": "Mar", "y": 180}, {"x": "Apr", "y": 220}]}
   ]}

kpi:
  {"type": "kpi", "value": "4,381", "label": "Active Copilot Users", "subtitle": "across 3 platforms"}

donut:
  {"type": "donut", "title": "Usage Tiers",
   "data": [{"label": "Light", "value": 144}, {"label": "Moderate", "value": 18}]}

funnel:
  {"type": "funnel", "data": [{"label": "Licensed", "value": 5000}, {"label": "Active", "value": 4381}]}

heatmap:
  {"type": "heatmap", "rows": ["Finance", "Legal"], "columns": ["Chat", "M365"],
   "values": [[5.5, 4.3], [3.0, 4.3]]}

table:
  {"type": "table", "columns": ["Department", "Sessions/wk"],
   "rows": [["Finance", "5.5"], ["Legal", "3.26"]], "highlight_col": 0}

gauge:
  {"type": "gauge", "value": "79", "max": 100, "label": "Return Rate"}

Output format:
{
  "deck_title": "Compelling story-driven title",
  "deck_subtitle": "Agents · Chat · M365 Copilot · Month Year – Month Year",
  "executive_summary": [
    "Finding with specific number -> business implication",
    ...  (5 bullets total, cross-page synthesis)
  ],
  "recommendations": [
    "Action: Specific recommendation with expected outcome",
    ...  (3-5 bullets)
  ],
  "slides": [
    {
      "slide_number": 1,
      "title": "Insight-led title",
      "headline": "Clear takeaway message",
      "insights": [
        {
          "text": "Bold punchy line || Supporting evidence with queried value [MeasureName]",
          "chart": {"type": "kpi", "value": "4,381", "label": "Active Users"}
        },
        {
          "text": "Bold punchy line || Pattern or opportunity from data",
          "chart": {"type": "bar", "title": "By Dept", "data": [...]}
        },
        {
          "text": "Bold punchy line || Actionable recommendation",
          "chart": null
        }
      ],
      "numbers_used": ["exact value from DAX result", "..."]
    }
  ]
}
""")
    print("-" * 70)


# ---------------------------------------------------------------------------
# Insight verification — catches missing charts & weak insights before build
# ---------------------------------------------------------------------------

def verify_insights(insights_file: str = "temp/insights.json",
                    request_file: str = "temp/analysis_request.json") -> dict:
    """Verify generated insights for common quality problems.

    Checks for:
      1. Slides where every chart spec is null (missing charts)
      2. Slide count mismatch (insights vs analysis request)
      3. Generic / vanilla headline patterns
      4. Missing executive summary or recommendations

    Returns a dict:
      {"passed": bool, "warnings": [str], "errors": [str]}
    """
    warnings: list = []
    errors: list = []

    # Load insights
    try:
        with open(insights_file, 'r', encoding='utf-8') as f:
            insights = json.load(f)
    except FileNotFoundError:
        return {"passed": False, "warnings": [],
                "errors": [f"Insights file not found: {insights_file}"]}
    except json.JSONDecodeError as e:
        return {"passed": False, "warnings": [],
                "errors": [f"Invalid JSON in {insights_file}: {e}"]}

    slides = insights.get("slides", [])
    if not slides:
        return {"passed": False, "warnings": [],
                "errors": ["No slides found in insights JSON"]}

    # ── Check 1: Missing charts ──────────────────────────────────────────
    for s in slides:
        slide_num = s.get("slide_number", "?")
        insight_items = s.get("insights", [])
        has_chart = False
        for item in insight_items:
            if isinstance(item, dict) and item.get("chart"):
                has_chart = True
                break

        if not has_chart:
            headline = s.get("headline", "")
            if headline.lower().strip() != "insufficient data for analysis":
                warnings.append(
                    f"Slide {slide_num}: No chart spec — will use screenshot fallback. "
                    f"Review whether chart data is extractable."
                )

    # ── Check 2: Slide count ─────────────────────────────────────────────
    try:
        with open(request_file, 'r', encoding='utf-8') as f:
            request = json.load(f)
        expected = request.get("total_slides", 0)
        actual = len(slides)
        if expected and actual < expected:
            warnings.append(
                f"Slide count mismatch: expected {expected}, got {actual}. "
                f"Some dashboard pages may be missing from insights."
            )
    except Exception:
        pass  # Non-critical — request file may not exist in --verify mode

    # ── Check 3: Generic / vanilla headlines ─────────────────────────────
    _VANILLA_PATTERNS = [
        "groups form",
        "org structure",
        "overview of",
        "shows the",
        "displays the",
        "presents the",
        "summary of",
        "breakdown of",
        "distribution of",
        "there are",
        "this page shows",
        "this slide shows",
    ]
    for s in slides:
        headline = (s.get("headline") or "").lower()
        for pat in _VANILLA_PATTERNS:
            if pat in headline:
                warnings.append(
                    f"Slide {s.get('slide_number', '?')}: Headline may be too "
                    f"generic (matched '{pat}'). Headlines should answer 'so what?' "
                    f"for an executive — not describe the chart."
                )
                break

    # ── Check 4: Missing executive summary / recommendations ─────────────
    if not insights.get("executive_summary"):
        errors.append("Missing 'executive_summary' — required for the summary slide.")
    if not insights.get("recommendations"):
        warnings.append("Missing 'recommendations' — strongly recommended.")
    if not insights.get("deck_title") or insights.get("deck_title") == "Executive Insights":
        warnings.append("Deck title is generic. Generate a compelling, specific title.")

    passed = len(errors) == 0

    # ── Print report ─────────────────────────────────────────────────────
    print("\n" + "=" * 70)
    print("INSIGHT VERIFICATION REPORT")
    print("=" * 70)

    if errors:
        for e in errors:
            print(f"  ERROR  {e}")
    if warnings:
        for w in warnings:
            print(f"  WARN   {w}")
    if not errors and not warnings:
        print("  OK  All checks passed — insights look good")

    print(f"\n  Result: {'PASS' if passed else 'FAIL'} "
          f"({len(errors)} error(s), {len(warnings)} warning(s))")
    print("=" * 70)

    return {"passed": passed, "warnings": warnings, "errors": errors}


def build_presentation_from_insights(source_path, output_path, insights_file, *, vector_charts=False):
    """Build final presentation using Claude's insights.

    Args:
        vector_charts: If True, use matplotlib vector charts instead of PBI screenshots.
                       Default is False (screenshots).
    """

    print("\n" + "=" * 70)
    print("BUILDING FINAL PRESENTATION")
    print("=" * 70)

    # ── Pre-build verification ──────────────────────────────────────────
    vresult = verify_insights(insights_file)
    if not vresult["passed"]:
        print("\n  Build aborted — fix the errors above before building.")
        return
    # Warnings are non-blocking: print them and continue

    # Load Claude's insights
    with open(insights_file, 'r', encoding='utf-8') as f:
        insights_data = json.load(f)

    print(f"\nOK Loaded insights for {len(insights_data.get('slides', []))} slides")

    # Use existing smart converter's rendering engine
    from lib.rendering.builder import render_presentation
    from lib.analysis.insights import Insight

    # Convert Claude's insights to expected format
    # Key by slide_number (not title) to ensure all slides are included
    from lib.analysis.insights import parse_bullet_points
    insights = {}
    for slide_insight in insights_data.get('slides', []):
        slide_num = slide_insight.get('slide_number')
        if slide_num:
            # Use slide_number as key for guaranteed matching
            insights[slide_num] = Insight(
                headline=slide_insight['headline'],
                bullet_points=parse_bullet_points(slide_insight.get('insights', [])),
                source_numbers=slide_insight.get('numbers_used', [])
            )
        else:
            # Fallback to title-based for backward compatibility
            title = slide_insight.get('title', '')
            if title:
                insights[title] = Insight(
                    headline=slide_insight['headline'],
                    bullet_points=parse_bullet_points(slide_insight.get('insights', [])),
                    source_numbers=slide_insight.get('numbers_used', [])
                )

    # Add executive summary, recommendations, deck title, and rendering mode as special keys
    insights['__vector_charts__'] = vector_charts
    if 'executive_summary' in insights_data:
        insights['__executive_summary__'] = insights_data['executive_summary']
    if 'recommendations' in insights_data:
        insights['__recommendations__'] = insights_data['recommendations']
    if 'deck_title' in insights_data:
        insights['__deck_title__'] = insights_data['deck_title']
    if 'deck_subtitle' in insights_data:
        insights['__deck_subtitle__'] = insights_data['deck_subtitle']

    # Render presentation
    vector_charts = insights_data.get('__vector_charts__', False)
    print(f"\nRendering presentation... (mode: {'vector charts' if vector_charts else 'PBI screenshots'})")
    render_presentation(source_path, insights, output_path, vector_charts=vector_charts)

    print(f"\nOK Created: {output_path}")

    # Validate
    from lib.rendering.validator import validate_output
    passed, report = validate_output(insights)
    print("\n" + report)

    # Clean up stale insight files so they don't bleed into the next conversion
    _cleanup_insight_files()

    print("\n" + "=" * 70)
    print("OK CONVERSION COMPLETE")
    print("=" * 70)


def _cleanup_insight_files():
    """Delete per-run insight artefacts after a successful build."""
    import os
    for fname in ("temp/insights.json", "temp/write_insights.py"):
        try:
            os.remove(fname)
        except FileNotFoundError:
            pass


def generate_output_filename(source_path):
    """Generate output filename from source (e.g., dashboard.pptx -> dashboard_executive.pptx)"""
    from pathlib import Path
    source = Path(source_path)
    # For directories (PBIP project folder), use the folder name
    if source.is_dir():
        # Look for a .pbip file to get the project name
        pbip_files = list(source.glob('*.pbip'))
        if pbip_files:
            return str(source / f"{pbip_files[0].stem}_executive.pptx")
        return str(source / f"{source.name}_executive.pptx")
    # Output is always .pptx regardless of input format (PPTX, PDF, or PBIP)
    return str(source.parent / f"{source.stem}_executive.pptx")


def main():
    parser = argparse.ArgumentParser(
        prog='convert_dashboard.py',
        description='Convert Power BI dashboards to executive presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Default (PBI screenshots):
  python convert_dashboard.py "dashboard.pbip"

  # Vector charts mode:
  python convert_dashboard.py "dashboard.pbip" --vector-charts

  # Custom output name:
  python convert_dashboard.py "dashboard.pptx" --output executive.pptx

  # Manual two-step workflow:
  python convert_dashboard.py "dashboard.pptx" --prepare
  python convert_dashboard.py --build --output executive.pptx
  python convert_dashboard.py --build --output executive.pptx --vector-charts
"""
    )

    parser.add_argument('source', nargs='?', default=None,
                       help='Source PowerPoint, PDF, or PBIP file/directory')
    parser.add_argument('--output', help='Output PowerPoint file (default: source_executive.pptx)')
    parser.add_argument('--prepare', action='store_true',
                       help='Prepare slides for analysis (Step 1 only)')
    parser.add_argument('--build', action='store_true',
                       help='Build final presentation from insights (Step 3 only)')
    parser.add_argument('--verify', action='store_true',
                       help='Verify insights JSON for missing charts, weak headlines, etc.')
    parser.add_argument('--insights', default='temp/insights.json',
                       help='Path to insights JSON (default: temp/insights.json)')
    parser.add_argument('--auto', action='store_true',
                       help='Auto mode: skip interactive prompt (for non-interactive environments)')
    parser.add_argument('--vector-charts', action='store_true',
                       help='Use matplotlib vector charts instead of PBI page screenshots (default: screenshots)')
    parser.add_argument('--context', default=None,
                       help='Optional analysis focus injected into the prompt, e.g. "spotlight Group A"')
    parser.add_argument('--assistant', default='auto', choices=['claude', 'copilot', 'auto'],
                       help='Which assistant to use for insights (claude, copilot, auto)')

    # Support legacy --source flag (maps to positional)
    # Parse known args first to handle --source gracefully
    args, remaining = parser.parse_known_args()
    if not args.source and remaining:
        # Check for legacy --source
        for i, arg in enumerate(remaining):
            if arg == '--source' and i + 1 < len(remaining):
                args.source = remaining[i + 1]
                break

    args = parser.parse_args()

    # ========================================================================
    # SINGLE-COMMAND WORKFLOW: Orchestrate all 3 steps automatically
    # ========================================================================
    if args.source and not args.prepare and not args.build:
        assistant = _resolve_assistant(args.assistant)
        # Auto-generate output filename if not provided
        output_path = args.output or generate_output_filename(args.source)
        print("\n" + "=" * 70)
        print("AUTOMATED CONVERSION WORKFLOW")
        print("=" * 70)
        print(f"\nSource: {args.source}")
        print(f"Output: {output_path}")
        print("\nThis will run all 3 steps automatically:")
        print("  1. Extract dashboard images")
        print("  2. Assistant analyzes and generates insights")
        print("  3. Build executive presentation")

        # STEP 1: Prepare slides for analysis
        print("\n" + "=" * 70)
        print("STEP 1: EXTRACTING DASHBOARDS")
        print("=" * 70)
        request_file = prepare_for_analysis(args.source, use_text_layer=(assistant == 'copilot'))

        # STEP 2: Trigger assistant analysis
        if assistant == 'copilot':
            show_copilot_instructions(request_file, context=args.context)
        else:
            trigger_claude_analysis(request_file, context=args.context)

        # Wait for Claude to generate insights
        # Auto-detect non-interactive environments and poll for insights file
        print("\n" + "=" * 70)
        print("Waiting for insights file...")
        print("=" * 70)

        max_wait = 300  # 5 minutes max
        wait_interval = 2  # Check every 2 seconds
        elapsed = 0

        while elapsed < max_wait:
            if Path(args.insights).exists():
                try:
                    # Verify it's valid JSON and has slides
                    with open(args.insights, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                        if 'slides' in data and len(data['slides']) > 0:
                            print(f"OK Claude analysis complete ({elapsed}s)")
                            break
                except (json.JSONDecodeError, KeyError):
                    pass  # File exists but not ready yet

            time.sleep(wait_interval)
            elapsed += wait_interval
        else:
            if assistant == 'copilot':
                print(f"Warning: Insights file not ready after {max_wait}s.")
                print("Run Copilot Chat to generate temp/insights.json and re-run this command.")
                return 0
            print(f"Warning: Insights file not ready after {max_wait}s, proceeding anyway...")

        print("=" * 70)

        # STEP 3: Build final presentation
        print("\n" + "=" * 70)
        print("STEP 3: BUILDING PRESENTATION")
        print("=" * 70)
        build_presentation_from_insights(args.source, output_path, args.insights,
                                       vector_charts=getattr(args, 'vector_charts', False))

        print("\n" + "=" * 70)
        print("CONVERSION COMPLETE!")
        print("=" * 70)
        print(f"\nExecutive presentation created: {output_path}")
        return 0

    # ========================================================================
    # MANUAL WORKFLOW: Individual steps
    # ========================================================================
    if args.verify:
        # Standalone verification (no build)
        result = verify_insights(args.insights)
        return 0 if result["passed"] else 1

    if args.prepare or (args.source and not args.build):
        # Step 1: Prepare slides
        if not args.source:
            print("Error: --source required")
            return 1

        assistant = _resolve_assistant(args.assistant)
        request_file = prepare_for_analysis(args.source, use_text_layer=(assistant == 'copilot'))
        if assistant == 'copilot':
            show_copilot_instructions(request_file, context=args.context)
        else:
            show_claude_instructions(request_file, context=args.context)

    elif args.build:
        # Step 3: Build final presentation
        # Get source from request file
        with open('temp/analysis_request.json', 'r') as f:
            request = json.load(f)
            source_path = request['source_file']

        # Auto-generate output filename if not provided
        output_path = args.output or generate_output_filename(source_path)

        build_presentation_from_insights(source_path, output_path, args.insights,
                                       vector_charts=getattr(args, 'vector_charts', False))

    else:
        parser.print_help()
        return 1

    return 0


if __name__ == '__main__':
    sys.exit(main())
