"""
Tests for executive summary and recommendations slides
"""
import pytest
from pathlib import Path
from pptx import Presentation
from lib.rendering.builder import SlideBuilder


def test_add_executive_summary_slide():
    """Test executive summary slide creation"""
    # Create builder with dummy source
    source_prs = Presentation()
    builder = SlideBuilder(source_prs)

    # Add executive summary slide
    summary_bullets = [
        "134 Agent users from 1,275 total (11%) → significant automation adoption opportunity",
        "HR Generalists (140 actions/user) are 3-4x above average → proven champions",
        "Finance achieves 217 prompts/user (5.7x average) → successful pattern to replicate",
        "Platform shows concentration in top 20% users → expand engagement baseline",
        "External collaboration drives 73% of assisted value → prioritize external workflows"
    ]

    slide = builder.add_executive_summary_slide(summary_bullets)

    # Verify slide was created
    assert slide is not None
    assert len(builder.prs.slides) == 1

    # Verify slide has text shapes
    text_shapes = [s for s in slide.shapes if hasattr(s, 'text')]
    assert len(text_shapes) >= 2  # Title + bullets

    # Verify title
    title_text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and 'Executive Summary' in shape.text:
            title_text = shape.text
            break
    assert title_text == "Executive Summary"

    print("✓ Executive summary slide created successfully")


def test_add_recommendations_slide():
    """Test recommendations slide creation"""
    # Create builder with dummy source
    source_prs = Presentation()
    builder = SlideBuilder(source_prs)

    # Add recommendations slide
    recommendations = [
        "Pilot agent training with HR Generalists (140 actions/user) to establish best practices for departments below 50 actions/user",
        "Expand PowerPoint Copilot access to Finance (217 prompts/user pattern) to replicate proven workflow",
        "Target top 20% users for champion program to drive baseline engagement from 38 to 50+ prompts/user"
    ]

    slide = builder.add_recommendations_slide(recommendations)

    # Verify slide was created
    assert slide is not None
    assert len(builder.prs.slides) == 1

    # Verify slide has text shapes
    text_shapes = [s for s in slide.shapes if hasattr(s, 'text')]
    assert len(text_shapes) >= 2  # Title + recommendations

    # Verify title
    title_text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and 'Next Steps' in shape.text:
            title_text = shape.text
            break
    assert title_text == "Next Steps & Recommendations"

    print("✓ Recommendations slide created successfully")


def test_full_presentation_with_executive_slides():
    """Test complete presentation with executive summary and recommendations"""
    # Create builder with dummy source
    source_prs = Presentation()
    builder = SlideBuilder(source_prs)

    # Add title slide
    builder.add_title_slide("Test Dashboard", "Test Subtitle")

    # Add executive summary
    summary = [
        "Finding 1 with data → implication",
        "Finding 2 with data → implication",
        "Finding 3 with data → implication",
        "Finding 4 with data → implication",
        "Finding 5 with data → implication"
    ]
    builder.add_executive_summary_slide(summary)

    # Add content slide
    builder.add_insight_slide(
        slide_number=1,
        headline="Test Headline",
        insights=["Insight 1", "Insight 2"],
        source_image=None
    )

    # Add recommendations
    recommendations = [
        "Recommendation 1 with specifics",
        "Recommendation 2 with specifics",
        "Recommendation 3 with specifics"
    ]
    builder.add_recommendations_slide(recommendations)

    # Verify slide order
    assert len(builder.prs.slides) == 4

    # Verify slide order by checking titles
    slide_texts = []
    for slide in builder.prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_texts.append(shape.text)
                break

    assert any('Test Dashboard' in text for text in slide_texts)  # Title slide
    assert any('Executive Summary' in text for text in slide_texts)  # Executive summary
    assert any('Test Headline' in text for text in slide_texts)  # Content slide
    assert any('Next Steps' in text for text in slide_texts)  # Recommendations

    print("✓ Full presentation structure validated")


def test_executive_summary_max_bullets():
    """Test that executive summary limits to 5 bullets"""
    source_prs = Presentation()
    builder = SlideBuilder(source_prs)

    # Try to add 7 bullets (should cap at 5)
    summary = [f"Bullet {i}" for i in range(7)]
    slide = builder.add_executive_summary_slide(summary)

    # Count bullets in slide
    bullet_count = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.text and '•' in paragraph.text:
                    bullet_count += 1

    assert bullet_count == 5, f"Expected 5 bullets, got {bullet_count}"

    print("✓ Executive summary correctly limits to 5 bullets")


if __name__ == '__main__':
    test_add_executive_summary_slide()
    test_add_recommendations_slide()
    test_full_presentation_with_executive_slides()
    test_executive_summary_max_bullets()
    print("\n✓ All executive slides tests passed!")
